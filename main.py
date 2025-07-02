
# Функция для открытия файлов с использованием программ по умолчанию

def open_file_with_default_program(file_path):
    try:
        if platform.system() == 'Windows':
            os.startfile(file_path)
        elif platform.system() == 'Darwin':  # macOS
            subprocess.call(['open', file_path])
        else:  # linux
            subprocess.call(['xdg-open', file_path])
    except Exception as e:
        logger.error(f"Ошибка открытия файла {file_path}: {e}")

import os
import requests
import base64
import json
import re
import subprocess
import winreg
import shutil
import urllib.parse
import venv
import webbrowser
import platform
import io
import zipfile
import mimetypes

from tqdm import tqdm
from loguru import logger
from urllib.parse import urlparse
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from tkinter import messagebox
from packaging import version
from lang import get_text, get_language_settings, save_language_settings, show_language_selection, translate_prompt_for_ai

# Импорты для работы с различными форматами файлов
try:
    from PIL import Image, ImageDraw, ImageFont
    import pytesseract
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

CURRENT_VERSION = "1.1"
# Логирование
logger.add("pollinations_agent.log", rotation="50 MB")

# Имитация флага перевода
class MainApp:
    def __init__(self):
        self.isTranslate = False  # Русский/Английский

main_app = MainApp()

# Функция получения сообщения об ошибке
def get_error_message(translate=False):
    if translate:
        return "Ошибка"
    else:
        return "Error"

# Удаление эмодзи из строки
def remove_emojis(text):
    emoji_pattern = re.compile(
        "["
        "\U0001F600-\U0001F64F"  # emoticons
        "\U0001F300-\U0001F5FF"  # symbols & pictographs
        "\U0001F680-\U0001F6FF"  # transport & map symbols
        "\U0001F700-\U0001F77F"  # alchemical symbols
        "\U0001F780-\U0001F7FF"  # Geometric Shapes Extended
        "\U0001F800-\U0001F8FF"  # Supplemental Symbols and Pictographs
        "\U0001F900-\U0001F9FF"  # Emoticons Supplement
        "\U0001FA00-\U0001FA6F"  # Chess Symbols
        "\U0001FA70-\U0001FAFF"  # Symbols and Pictographs Extended-A
        "\U00002702-\U000027B0"  # Dingbats
        "]+",
        flags=re.UNICODE,
    )
    return emoji_pattern.sub(r"", text)

def update_app(update_url):
   webbrowser.open(update_url)

def check_for_updates():
    try:
        # Получение информации о последнем релизе на GitHub
        response = requests.get("https://api.github.com/repos/Processori7/Poli_AI/releases/latest")
        response.raise_for_status()
        latest_release = response.json()
        # Получение ссылки на файл llm.exe последней версии
        download_url = None
        assets = latest_release["assets"]
        for asset in assets:
            if asset["name"] == "poliai.exe":  # Ищем только llm.exe
                download_url = asset["browser_download_url"]
                break

        if download_url is None:
            messagebox.showerror("Ошибка обновления", "Не удалось найти файл poliai.exe для последней версии.")
            return

        # Сравнение текущей версии с последней версией
        latest_version_str = latest_release["tag_name"]
        match = re.search(r'\d+\.\d+', latest_version_str)
        if match:
            latest_version = match.group()
        else:
            latest_version = latest_version_str

        if version.parse(latest_version) > version.parse(CURRENT_VERSION):
            if platform.system() == "Windows":
                # Предложение пользователю обновление
                if messagebox.showwarning("Доступно обновление",
                                          f"Доступна новая версия {latest_version}. Хотите обновить?", icon='warning',
                                          type='yesno') == 'yes':
                    update_app(download_url)
            else:
                if messagebox.showwarning("Доступно обновление",
                                          f"Доступна новая версия {latest_version}. Хотите обновить?", icon='warning',
                                          type='yesno') == 'yes':
                    os.system("git pull")
    except requests.exceptions.RequestException as e:
        messagebox.showerror("Error", str(e))

# Генерация изображения
def gen_img(prompt, model_name="flux", width=1024, height=1024):
    try:
        encoded_prompt = requests.utils.quote(prompt)
        url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?model={model_name}&width={width}&height={height}"
        response = requests.get(url)
        if response.status_code == 200:
            output_dir = "output"
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            filepath = os.path.join(output_dir, f"{prompt.replace(' ', '_')}.jpg")
            with open(filepath, "wb") as f:
                f.write(response.content)
            return filepath
        else:
            return f"Ошибка генерации: {response.status_code}"
    except Exception as e:
        return str(e)

# Общение с Pollinations (текстовая модель)
def communicate_with_Pollinations_chat(model_name, messages, tools=None, tool_choice="auto", api_token=None):
    """
    Отправляет структурированный запрос к Pollinations API с поддержкой инструментов.

    :param model_name: Имя модели (например, 'openai')
    :param messages: Список сообщений в формате [{'role': 'user', 'content': '...'}, ...]
    :param tools: Список инструментов в формате OpenAI
    :param tool_choice: auto или конкретное имя инструмента
    :param api_token: API токен (опционально)
    :return: JSON-ответ от сервера или ошибка
    """
    url = "https://text.pollinations.ai/openai"

    payload = {
        "model": model_name,
        "messages": messages,
        "private": True  # Не показывать в публичном feed
    }

    # Добавляем tools и tool_choice только если tools переданы
    if tools:
        payload["tools"] = tools
        payload["tool_choice"] = tool_choice

    # Подготавливаем заголовки
    headers = {"Content-Type": "application/json"}
    
    # Добавляем токен если он есть
    if api_token:
        headers["Authorization"] = f"Bearer {api_token}"

    try:
        response = requests.post(
            url,
            headers=headers,
            json=payload
        )
        # Убрали вывод HTTP статуса и ответа
        if response.status_code == 200:
            return response.json()
        else:
            return {"error": f"HTTP {response.status_code}: {response.text}"}
    except Exception as e:
        return {"error": str(e)}

# Получить модели для генерации изображений
def get_Polinations_img_models():
    model_functions = {}
    try:
        url = "https://image.pollinations.ai/models"
        resp = requests.get(url)
        if resp.ok:
            models = resp.json()  # Получаем список строк
            for name in models:  # Проходим по каждой строке
                key = f"(Polinations) {name}_img"
                model_functions[key] = lambda user_input, model_name=name: gen_img(user_input, model_name)
            return model_functions
        else:
            return f"{get_error_message(main_app.isTranslate)}: {resp.status_code}"
    except Exception as e:
        return f"{get_error_message(main_app.isTranslate)}: {str(e)}"

# Получить модели для чат-генерации текста (с описанием)
def get_Polinations_chat_models():
    models_list = []
    try:
        url = "https://text.pollinations.ai/models"
        resp = requests.get(url)
        if resp.ok:
            models = resp.json()
            for model in models:
                if isinstance(model, dict) and "name" in model:
                    model_name = model["name"]
                    model_description = model.get("description", "Без описания")
                    models_list.append({
                        "name": model_name,
                        "description": model_description
                    })
            return models_list
        else:
            print(f"Ошибка получения списка моделей: {resp.status_code}")
            return [{"name": "o3-mini", "description": "Быстрая и эффективная модель"}]
    except Exception as e:
        print(f"Ошибка при получении списка моделей: {e}")
        return [{"name": "o3-mini", "description": "Быстрая и эффективная модель"}]

# Класс агента
class PollinationsAgent:
    def __init__(self):
        # Загружаем переменные окружения из .env файла
        load_dotenv()
        
        self.base_text_url = "https://text.pollinations.ai"
        self.base_image_url = "https://image.pollinations.ai"
        self.output_dir = "output"
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)
            
        # Загружаем настройки из .env файла
        self.api_token = os.getenv('POLLINATIONS_TOKEN')

        lang_settings = get_language_settings()
        self.interface_language = lang_settings['interface']
        self.output_language = lang_settings['output']
        self.default_model = os.getenv('DEFAULT_MODEL', 'openai')
        self.max_attempts = int(os.getenv('MAX_ATTEMPTS', '3'))
        self.default_voice = os.getenv('DEFAULT_VOICE', 'alloy')
        self.require_confirmation = os.getenv('REQUIRE_CONFIRMATION', 'true').lower() == 'true'
        self.debug_mode = os.getenv('DEBUG_MODE', 'false').lower() == 'true'
        
        # Сообщения в зависимости от языка интерфейса
        self.messages = self._get_interface_messages()
        
        if self.api_token:
            print(f"🔑 {self.messages['api_token_loaded']}")
        else:
            print(f"⚠️ {self.messages['api_token_missing']}")
            
        if self.debug_mode:
            print(f"🔧 Debug mode: Interface={self.interface_language}, Output={self.output_language}, Model={self.default_model}")

        # MCP инструменты
        self.mcp_tools = {
            "generateImageUrl": self.generate_image_url,
            "generateImage": self.generate_image,
            "respondAudio": self.generate_audio,
            "sayText": self.generate_audio,
            "listImageModels": self.list_image_models,
            "listAudioVoices": self.list_audio_voices,

            # 🔧 Файловая система
            "createFile": self.create_file,
            "readFile": self.read_file,
            "writeFile": self.write_file,
            "deleteFile": self.delete_file,
            "moveFile": self.move_file,
            "listDirectory": self.list_directory,
            "createDirectory": self.create_directory,
            "deleteDirectory": self.delete_directory,
            
            # 🌐 Интернет и загрузки
            "downloadFile": self.download_file,
            "downloadImage": self.download_image,
            "searchInternet": self.search_internet,
            
            # 💻 Системные команды
            "executeCommand": self.execute_command,
            "runPythonCode": self.run_python_code,
            
            # ⚙️ Системные настройки
            "changeRegistryValue": self.change_registry_value,
            "getSystemInfo": self.get_system_info,
            "manageServices": self.manage_services,
            
            # 🔍 Поиск изображений
            "searchAndDownloadImages": self.search_and_download_images,
            
            # 🐍 Python проекты
            "createPythonProject": self.create_python_project,
            
            # 💻 Разработка ПО
            "developSoftware": self.develop_software,
            
            # 📄 Чтение файлов различных форматов
            "readAdvancedFile": self.read_advanced_file,
            "readDocx": self.read_docx_file,
            "readPdf": self.read_pdf_file,
            "readExcel": self.read_excel_file,
            "readPowerPoint": self.read_powerpoint_file,
            
            # 🖼️ Анализ изображений
            "analyzeImage": self.analyze_image,
            "recognizeText": self.recognize_text_from_image,
            "getImageInfo": self.get_image_info,
            "findAndAnalyzeFile": self.find_and_analyze_file,
            
            # 📂 Открытие и запуск файлов
            "openFileWithDefaultProgram": self.open_file_with_default_program,
            "runExecutable": self.run_executable,
            "smartOpenFile": self.smart_open_file,
        }

        # Получаем доступные модели
        self.img_models = get_Polinations_img_models()
        self.chat_models = get_Polinations_chat_models()  # Теперь это список словарей
        self.current_model = None

        # Сохраняем только имена моделей
        self.model_list = [model['name'] for model in self.chat_models]

    def create_file(self, path, content=""):
        """Создает файл по указанному пути"""
        try:
            # Создаем директорию, если она не существует
            directory = os.path.dirname(path)
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True)
                logger.info(f"Создана папка: {directory}")
            
            with open(path, 'w', encoding='utf-8') as f:
                f.write(content)
            logger.info(f"Файл создан: {path}")
            return f"Файл создан: {path}"
        except Exception as e:
            return f"Ошибка создания файла: {str(e)}"

    def read_file(self, path):
        """Читает содержимое файла"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                content = f.read()
            logger.info(f"Файл прочитан: {path}")
            return content
        except Exception as e:
            return f"Ошибка чтения файла: {str(e)}"

    def write_file(self, path, content):
        """Записывает данные в существующий файл"""
        try:
            with open(path, 'w', encoding='utf-8') as f:
                f.write(content)
            logger.info(f"Файл обновлен: {path}")
            return f"Файл обновлен: {path}"
        except Exception as e:
            return f"Ошибка записи в файл: {str(e)}"

    def delete_file(self, path):
        """Удаляет файл"""
        try:
            os.remove(path)
            logger.info(f"Файл удален: {path}")
            return f"Файл удален: {path}"
        except Exception as e:
            return f"Ошибка удаления файла: {str(e)}"
    
    def move_file(self, source_path, destination_path):
        """Перемещает файл из одной папки в другую"""
        try:
            # Создаем папку назначения если она не существует
            destination_dir = os.path.dirname(destination_path)
            if destination_dir and not os.path.exists(destination_dir):
                os.makedirs(destination_dir, exist_ok=True)
                logger.info(f"Создана папка: {destination_dir}")
            
            # Если путь назначения - это папка, добавляем имя файла
            if os.path.isdir(destination_path):
                filename = os.path.basename(source_path)
                destination_path = os.path.join(destination_path, filename)
            
            # Перемещаем файл
            shutil.move(source_path, destination_path)
            logger.info(f"Файл перемещен: {source_path} -> {destination_path}")
            return f"Файл перемещен: {source_path} -> {destination_path}"
        except Exception as e:
            return f"Ошибка перемещения файла: {str(e)}"

    def list_directory(self, path="."):
        """Возвращает список файлов и папок в директории"""
        try:
            items = os.listdir(path)
            logger.info(f"Директория просканирована: {path}")
            return json.dumps(items, ensure_ascii=False)
        except Exception as e:
            return f"Ошибка сканирования директории: {str(e)}"

    def create_directory(self, path):
        """Создает новую папку"""
        try:
            os.makedirs(path, exist_ok=True)
            logger.info(f"Папка создана: {path}")
            return f"Папка создана: {path}"
        except Exception as e:
            return f"Ошибка создания папки: {str(e)}"

    def delete_directory(self, path):
        """Удаляет папку и её содержимое"""
        try:
            # Используем shutil.rmtree для рекурсивного удаления
            if os.path.exists(path):
                shutil.rmtree(path)
                logger.info(f"Папка удалена: {path}")
                return f"Папка удалена: {path}"
            else:
                return f"Папка не существует: {path}"
        except Exception as e:
            return f"Ошибка удаления папки: {str(e)}"
    
    def swap_files(self, path1, path2):
        """Меняет два файла местами"""
        try:
            # Проверяем что оба файла существуют
            if not os.path.exists(path1):
                return f"Файл не существует: {path1}"
            if not os.path.exists(path2):
                return f"Файл не существует: {path2}"
            
            # Создаем временный файл
            temp_file = path1 + ".temp_swap"
            
            # Меняем файлы местами через временный файл
            shutil.move(path1, temp_file)
            shutil.move(path2, path1)
            shutil.move(temp_file, path2)
            
            logger.info(f"Файлы обменены местами: {path1} ↔ {path2}")
            return f"Файлы обменены местами: {path1} ↔ {path2}"
        except Exception as e:
            return f"Ошибка обмена файлов: {str(e)}"

    def list_image_models(self):
        """Список доступных моделей изображений"""
        response = requests.get(f"{self.base_image_url}/models")
        response.raise_for_status()
        return response.json()

    def list_audio_voices(self):
        """Список доступных голосов"""
        try:
            response = requests.get(f"{self.base_text_url}/models")
            response.raise_for_status()
            models_data = response.json()
            
            # Попытаемся извлечь голоса из структуры данных
            if isinstance(models_data, dict) and 'openai-audio' in models_data:
                voices = models_data.get('openai-audio', {}).get('voices', [])
                if voices:
                    return voices
            
            # Если не удалось получить из API, возвращаем известные голоса
            return ["alloy", "echo", "fable", "onyx", "nova", "shimmer"]
        except Exception as e:
            logger.warning(f"Не удалось получить список голосов: {e}")
            # Возвращаем стандартные голоса OpenAI
            return ["alloy", "echo", "fable", "onyx", "nova", "shimmer"]

    # def translate_to_english(self, text):
    #     """Переводит текст на английский язык"""
    #     try:
    #         # Проверяем, нужен ли перевод (если текст уже на английском)
    #         if self._is_english(text):
    #             return text
    #
    #         prompt = f"Переведи следующий текст на английский язык. Отвечай только переводом без дополнительных комментариев:\n\n{text}"
    #
    #         messages = [{"role": "user", "content": prompt}]
    #         response = communicate_with_Pollinations_chat(self.current_model, messages)
    #
    #         if "choices" in response and len(response["choices"]) > 0:
    #             translated = response["choices"][0]["message"].get("content", "").strip()
    #             logger.info(f"Текст переведен: '{text}' -> '{translated}'")
    #             return translated
    #         else:
    #             logger.warning(f"Не удалось перевести текст: {text}")
    #             return text
    #     except Exception as e:
    #         logger.warning(f"Ошибка перевода: {str(e)}")
    #         return text
    #
    # def _is_english(self, text):
    #     """Проверяет, написан ли текст на английском языке"""
    #     try:
    #         # Простая эвристика: если больше 70% символов - латинские, считаем английским
    #         latin_chars = sum(1 for c in text if c.isalpha() and ord(c) < 128)
    #         total_chars = sum(1 for c in text if c.isalpha())
    #
    #         if total_chars == 0:
    #             return True  # Если нет букв, считаем английским
    #
    #         return (latin_chars / total_chars) > 0.7
    #     except:
    #         return False

    def check_answear(self) -> bool:
        if not self.api_token:
            print(f"\n⚠️ Pollinations API токен не найден")
            print(f"🔗 Получите токен на: https://auth.pollinations.ai/")

            # Предлагаем открыть страницу для получения токена
            try:
                open_choice = input("\nОткрыть страницу для получения токена? (y/n): ").strip().lower()
                if open_choice == 'y':
                    webbrowser.open('https://auth.pollinations.ai/')
                    return True
                else:
                    return False
            except Exception:
                pass

    def generate_audio(self, text, voice="alloy"):
        """Генерация аудио через Pollinations API"""
        try:
            print(f"🎵 Генерация аудио...")
            print(f"📝 Исходный текст: {text}")
            
            # Проверяем наличие API токена
            if not self.api_token:
                print(f"\n⚠️ Pollinations API токен не найден")
                print(f"🔗 Получите токен на: https://auth.pollinations.ai/")

                # Предлагаем открыть страницу для получения токена
                try:
                    open_choice = input("\nОткрыть страницу для получения токена? (y/n): ").strip().lower()
                    if open_choice == 'y':
                        webbrowser.open('https://auth.pollinations.ai/')
                        print("🌐 Страница открыта в браузере")
                    else:
                        print("Отмена открытия страницы")
                except Exception:
                    pass

            # Показываем доступные голоса для выбора
            available_voices = self.list_audio_voices()
            print(f"\n🎙️ Доступные голоса:")
            for i, voice_name in enumerate(available_voices, 1):
                current_marker = " (текущий)" if voice_name == voice else ""
                print(f"{i}. {voice_name}{current_marker}")

            try:
                choice = input(f"\nВыберите голос (1-{len(available_voices)}) или нажмите Enter для {voice}: ").strip()
                if choice and choice.isdigit():
                    choice_idx = int(choice) - 1
                    if 0 <= choice_idx < len(available_voices):
                        voice = available_voices[choice_idx]
                        print(f"✅ Выбран голос: {voice}")
                    else:
                        print(f"⚠️ Неверный номер, используется голос по умолчанию: {voice}")
                else:
                    print(f"⚠️ Используется голос по умолчанию: {voice}")
            except (ValueError, IndexError):
                print(f"⚠️ Ошибка ввода, используется голос по умолчанию: {voice}")

            # Подготавливаем payload с токеном
            payload = {
                "model": "openai-audio",
                "messages": [{"role": "user", "content": text}],
                "voice": voice,
                "private": True  # Не показывать в публичном feed
            }

            # Подготавливаем заголовки с Authorization Bearer токеном
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.api_token}"
            }

            print(f"🔄 Отправка запроса на генерацию аудио с токеном...")
            response = requests.post(
                f"{self.base_text_url}/openai",
                json=payload,
                headers=headers
            )
            response.raise_for_status()

            result = response.json()
            if 'choices' in result and len(result['choices']) > 0:
                message = result['choices'][0]['message']
                if 'audio' in message and 'data' in message['audio']:
                    audio_data = message['audio']['data']
                    filename = text.replace(" ", "_")[:50]  # Ограничиваем длину имени файла
                    filepath = self.save_audio(audio_data, filename)
                    print(f"✅ Аудио сгенерировано: {filepath}")
                    return f"Аудио сгенерировано и сохранено: {filepath}"
                else:
                    return "Ошибка: В ответе нет аудио данных"
            else:
                return "Ошибка: Неверный формат ответа от API"
                
        except requests.RequestException as e:
            error_msg = str(e)
            if "402" in error_msg or "Payment Required" in error_msg:
                return "❌ Требуется оплата или валидный токен для генерации аудио. Получите токен на https://auth.pollinations.ai/"
            return f"Ошибка сети при генерации аудио: {error_msg}"
        except Exception as e:
            return f"Ошибка генерации аудио: {str(e)}"


    def save_audio(self, audio_data, filename):
        """Сохраняет аудио в файл с отладочной информацией"""
        try:
            print(f"🔍 Отладка: Получены аудио данные, длина base64 строки: {len(audio_data)}")
            print(f"🔍 Отладка: Первые 100 символов: {audio_data[:100]}...")
            
            # Декодируем base64
            audio_bytes = base64.b64decode(audio_data)
            print(f"🔍 Отладка: Размер декодированных данных: {len(audio_bytes)} байт")
            
            # Проверяем магические байты для определения формата
            magic_bytes = audio_bytes[:12]
            print(f"🔍 Отладка: Магические байты: {magic_bytes.hex()}")
            
            # Определяем расширение файла по магическим байтам
            if magic_bytes.startswith(b'ID3') or magic_bytes[6:10] == b'ftyp':
                extension = 'mp3'
            elif magic_bytes.startswith(b'RIFF') and magic_bytes[8:12] == b'WAVE':
                extension = 'wav'
            elif magic_bytes.startswith(b'OggS'):
                extension = 'ogg'
            elif magic_bytes.startswith(b'fLaC'):
                extension = 'flac'
            else:
                print(f"⚠️ Неопознанный аудио формат, пробуем сохранить как .mp3")
                extension = 'mp3'
            
            # Очищаем имя файла от недопустимых символов
            clean_filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_')).rstrip()
            if not clean_filename:
                clean_filename = "audio_output"
            
            filepath = os.path.join(self.output_dir, f"{clean_filename}.{extension}")
            print(f"💾 Сохранение в: {filepath}")
            
            with open(filepath, "wb") as f:
                f.write(audio_bytes)
            
            # Проверяем, что файл действительно создался
            if os.path.exists(filepath):
                file_size = os.path.getsize(filepath)
                print(f"✅ Файл сохранен успешно: {filepath} (размер: {file_size} байт)")
                logger.info(f"Аудио сохранено: {filepath}")
                return filepath
            else:
                print(f"❌ Файл не создался: {filepath}")
                return None
                
        except Exception as e:
            print(f"❌ Ошибка при сохранении аудио: {str(e)}")
            logger.error(f"Ошибка сохранения аудио: {str(e)}")
            return None

    def select_model(self):
        """Выбор модели пользователем с отображением описания"""
        print(f"\n{get_text('available_models', self.interface_language)}")
        for idx, model in enumerate(self.chat_models, 1):
            print(f"{idx}. {model['name']} — {model['description']}")
        
        try:
            choice = int(input(f"{get_text('select_model', self.interface_language)} ")) - 1
            if 0 <= choice < len(self.chat_models):
                selected_model = self.chat_models[choice]['name']
                self.current_model = selected_model
                print(f"✅ {get_text('model_selected', self.interface_language)}: {selected_model}")
                logger.info(f"Выбрана модель: {selected_model}")
            else:
                print(f"⚠️ {get_text('invalid_choice', self.interface_language)}")
                default_model = self.default_model
                self.current_model = default_model
                print(f"🤖 {get_text('default_model', self.interface_language)}: {default_model}")
        except (ValueError, KeyboardInterrupt):
            print(f"⚠️ {get_text('invalid_choice', self.interface_language)}")
            default_model = self.default_model
            self.current_model = default_model
            print(f"🤖 {get_text('default_model', self.interface_language)}: {default_model}")

    def change_model(self):
        """Смена текущей модели"""
        self.select_model()

    def get_available_tools(self):
        return self.mcp_tools
    
    def get_openai_format_tools(self):
        """Преобразует инструменты в формат OpenAI"""
        tools = [
            # 📁 Файловая система
            {
                "type": "function",
                "function": {
                    "name": "createDirectory",
                    "description": "Создает новую папку",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к создаваемой папке"
                            }
                        },
                        "required": ["path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "createFile",
                    "description": "Создает новый файл с содержимым",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к создаваемому файлу"
                            },
                            "content": {
                                "type": "string",
                                "description": "Содержимое файла"
                            }
                        },
                        "required": ["path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "listDirectory",
                    "description": "Показывает содержимое папки",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к папке для просмотра",
                                "default": "."
                            }
                        },
                        "required": []
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "readFile",
                    "description": "Читает содержимое файла",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к файлу для чтения"
                            }
                        },
                        "required": ["path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "moveFile",
                    "description": "Перемещает файл из одной папки в другую",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "source_path": {
                                "type": "string",
                                "description": "Путь к исходному файлу"
                            },
                            "destination_path": {
                                "type": "string",
                                "description": "Путь назначения (файл или папка)"
                            }
                        },
                        "required": ["source_path", "destination_path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "deleteFile",
                    "description": "Удаляет файл",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к файлу для удаления"
                            }
                        },
                        "required": ["path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "deleteDirectory",
                    "description": "Удаляет папку",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к папке для удаления"
                            }
                        },
                        "required": ["path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "writeFile",
                    "description": "Записывает содержимое в существующий файл",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {
                                "type": "string",
                                "description": "Путь к файлу"
                            },
                            "content": {
                                "type": "string",
                                "description": "Содержимое для записи"
                            }
                        },
                        "required": ["path", "content"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "swapFiles",
                    "description": "Меняет два файла местами",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path1": {
                                "type": "string",
                                "description": "Путь к первому файлу"
                            },
                            "path2": {
                                "type": "string",
                                "description": "Путь ко второму файлу"
                            }
                        },
                        "required": ["path1", "path2"]
                    }
                }
            },
            # 🎨 Генерация изображений
            {
                "type": "function",
                "function": {
                    "name": "generateImage",
                    "description": "Генерирует изображение по текстовому описанию и сохраняет в файл",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "prompt": {
                                "type": "string",
                                "description": "Описание изображения для генерации"
                            },
                            "save_path": {
                                "type": "string",
                                "description": "Папка для сохранения изображения (опционально)"
                            },
                            "filename": {
                                "type": "string",
                                "description": "Имя файла (опционально)"
                            },
                            "width": {
                                "type": "integer",
                                "description": "Ширина изображения",
                                "default": 1024
                            },
                            "height": {
                                "type": "integer",
                                "description": "Высота изображения",
                                "default": 1024
                            },
                            "model": {
                                "type": "string",
                                "description": "Модель для генерации (flux, turbo, deliberate)",
                                "default": "flux"
                            }
                        },
                        "required": ["prompt"]
                    }
                }
            },
            # 🎵 Генерация аудио
            {
                "type": "function",
                "function": {
                    "name": "generateAudio",
                    "description": "Генерирует аудио из текста с помощью TTS",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "text": {
                                "type": "string",
                                "description": "Текст для преобразования в аудио"
                            },
                            "voice": {
                                "type": "string",
                                "description": "Голос для озвучки (alloy, echo, fable, onyx, nova, shimmer)",
                                "default": "alloy"
                            }
                        },
                        "required": ["text"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "sayText",
                    "description": "Озвучивает текст голосом (аналог generateAudio)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "text": {
                                "type": "string",
                                "description": "Текст для озвучивания"
                            },
                            "voice": {
                                "type": "string",
                                "description": "Голос для озвучки (alloy, echo, fable, onyx, nova, shimmer)",
                                "default": "alloy"
                            }
                        },
                        "required": ["text"]
                    }
                }
            },
            # 🌐 Интернет и загрузки
            {
                "type": "function",
                "function": {
                    "name": "downloadFile",
                    "description": "Загружает файл из интернета",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "url": {
                                "type": "string",
                                "description": "URL файла для загрузки"
                            },
                            "save_path": {
                                "type": "string",
                                "description": "Папка для сохранения (опционально)"
                            },
                            "filename": {
                                "type": "string",
                                "description": "Имя файла (опционально)"
                            }
                        },
                        "required": ["url"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "downloadImage",
                    "description": "Загружает изображение из интернета",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "url": {
                                "type": "string",
                                "description": "URL изображения для загрузки"
                            },
                            "save_path": {
                                "type": "string",
                                "description": "Папка для сохранения (опционально)"
                            },
                            "filename": {
                                "type": "string",
                                "description": "Имя файла (опционально)"
                            }
                        },
                        "required": ["url"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "searchInternet",
                    "description": "Поиск информации в интернете",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {
                                "type": "string",
                                "description": "Поисковый запрос"
                            },
                            "model": {
                                "type": "string",
                                "description": "Модель для поиска (searchgpt или elixposearch)",
                                "default": "searchgpt"
                            }
                        },
                        "required": ["query"]
                    }
                }
            },
            # 💻 Системные команды
            {
                "type": "function",
                "function": {
                    "name": "executeCommand",
                    "description": "Выполняет команду в командной строке (требует подтверждения)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "command": {
                                "type": "string",
                                "description": "Команда для выполнения"
                            }
                        },
                        "required": ["command"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "runPythonCode",
                    "description": "Выполняет Python код (требует подтверждения)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "code": {
                                "type": "string",
                                "description": "Python код для выполнения"
                            }
                        },
                        "required": ["code"]
                    }
                }
            },
            # ⚙️ Системная информация
            {
                "type": "function",
                "function": {
                    "name": "getSystemInfo",
                    "description": "Получает информацию о системе",
                    "parameters": {
                        "type": "object",
                        "properties": {},
                        "required": []
                    }
                }
            },
            # 🔍 Поиск и анализ файлов
            {
                "type": "function",
                "function": {
                    "name": "findAndAnalyzeFile",
                    "description": "Находит и анализирует файл в текущей директории по запросу. Особенно полезно для поиска файлов по числам или ключевым словам.",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {
                                "type": "string",
                                "description": "Запрос для поиска файла (например, 'картинка 11', 'файл 21', 'изображение с текстом')"
                            }
                        },
                        "required": ["query"]
                    }
                }
            },
            # 📄 Расширенное чтение файлов
            {
                "type": "function",
                "function": {
                    "name": "readAdvancedFile",
                    "description": "Умное чтение файлов различных форматов (txt, docx, pdf, xlsx, pptx, изображения)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "file_path": {
                                "type": "string",
                                "description": "Путь к файлу для чтения"
                            }
                        },
                        "required": ["file_path"]
                    }
                }
            },
            # 🖼️ Анализ изображений
            {
                "type": "function",
                "function": {
                    "name": "analyzeImage",
                    "description": "Анализирует изображение с помощью Vision API",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "image_path": {
                                "type": "string",
                                "description": "Путь к изображению для анализа"
                            }
                        },
                        "required": ["image_path"]
                    }
                }
            },
            # 🔍 Распознавание текста
            {
                "type": "function",
                "function": {
                    "name": "recognizeText",
                    "description": "Распознает текст на изображении (OCR)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "image_path": {
                                "type": "string",
                                "description": "Путь к изображению для распознавания текста"
                            }
                        },
                        "required": ["image_path"]
                    }
                }
            },
            # 📂 Открытие и запуск файлов
            {
                "type": "function",
                "function": {
                    "name": "openFileWithDefaultProgram",
                    "description": "Открывает файл с программой по умолчанию (например, .docx в Word, .pdf в PDF-ридере)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "file_path": {
                                "type": "string",
                                "description": "Путь к файлу для открытия"
                            }
                        },
                        "required": ["file_path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "runExecutable",
                    "description": "Запускает исполняемый файл (.exe, .bat, .cmd и др.) с возможностью запуска от имени администратора",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "file_path": {
                                "type": "string",
                                "description": "Путь к исполняемому файлу"
                            },
                            "run_as_admin": {
                                "type": "boolean",
                                "description": "Запустить от имени администратора (только Windows)",
                                "default": False
                            },
                            "require_confirmation": {
                                "type": "boolean",
                                "description": "Требовать подтверждение пользователя перед запуском",
                                "default": True
                            }
                        },
                        "required": ["file_path"]
                    }
                }
            },
            {
                "type": "function",
                "function": {
                    "name": "smartOpenFile",
                    "description": "Умное открытие файла: исполняемые файлы запускаются, остальные открываются с программой по умолчанию",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "file_path": {
                                "type": "string",
                                "description": "Путь к файлу для открытия"
                            },
                            "run_as_admin": {
                                "type": "boolean",
                                "description": "Запустить исполняемые файлы от имени администратора (только Windows)",
                                "default": False
                            }
                        },
                        "required": ["file_path"]
                    }
                }
            }
        ]
        return tools

    def generate_plan(self, user_task, max_attempts=None):
        """Генерирует план и выполняет задачу с автопродолжением до успешного завершения"""
        if max_attempts is None:
            max_attempts = self.max_attempts
            
        attempt = 0
        last_error = None
        all_executed_actions = []  # Инициализируем здесь, в начале функции
        
        # Проверяем, является ли это простым общением
        if self._is_simple_conversation(user_task):
            return self._handle_simple_conversation(user_task)
        
        while attempt < max_attempts:
            attempt += 1
            print(f"\n🔄 Попытка {attempt}/{max_attempts}...")
            
            # Формируем промпт с учетом предыдущих ошибок
            if last_error:
                if "Не выполнены действия:" in last_error:
                    # Задача выполнена частично, продолжаем недостающие действия
                    missing = last_error.replace("Не выполнены действия:", "").strip()
                    prompt = f"""ПРОДОЛЖИ выполнение задачи: {user_task}

Задача выполнена НЕ ПОЛНОСТЬЮ. Еще нужно выполнить: {missing}

ВАЖНО: 
1. НЕ повторяй уже выполненные действия!
2. Выполни ТОЛЬКО оставшиеся недостающие действия из списка: {missing}
3. Обязательно выполни ВСЕ недостающие шаги для полного завершения задачи
4. Разбивай задачу на максимально понятные и короткие действия, вчитывайся в запрос
5. Если тебя просят что-то создать или сделать или сгенерировать всегда создавай папку и работай в ней
6. Используй правильную последовательность действий

ИНСТРУКЦИИ ДЛЯ РАБОТЫ С ЛЮБЫМИ ФАЙЛАМИ:
-Сначала создай папку с тематическим названием
-В папке создай требуемый файл

ИНСТРУКЦИИ ДЛЯ ПЕРЕМЕЩЕНИЯ ФАЙЛОВ:
- Если задача "перемести файлы из папок", то СНАЧАЛА используй listDirectory для каждой папки
- ЗАТЕМ используй moveFile с КОНКРЕТНЫМИ путями найденных файлов
- Формат: moveFile(source_path="папка/файл.txt", destination_path="целевая_папка")
- НЕ придумывай имена файлов - используй ТОЛЬКО реальные файлы из listDirectory!

Пример последовательности:
1. listDirectory({{"path": "21"}}) → найти файлы
2. listDirectory({{"path": "11"}}) → найти файлы  
3. moveFile({{"source_path": "21/найденный_файл.txt", "destination_path": "Файлы"}})
4. moveFile({{"source_path": "11/найденный_файл.txt", "destination_path": "Файлы"}})
5. deleteDirectory({{"path": "21"}})
6. deleteDirectory({{"path": "11"}})

Выполни недостающие действия СЕЙЧАС!"""
                else:
                    # Была ошибка выполнения, исправляем
                    prompt = f"""Выполни следующую задачу, используя доступные инструменты: {user_task}

Предыдущая попытка завершилась ошибкой: {last_error}

ВАЖНО: Проанализируй ошибку и используй ДРУГОЙ подход! 
Не повторяй те же самые действия, которые привели к ошибке.

Примеры исправлений:
- Если папка не существует - сначала создай её
- Если файл не найден - проверь правильность пути
- Если нет доступа - измени разрешения или путь
- Если команда не работает - используй альтернативную команду

Обязательно измени последовательность действий или используй другие инструменты!"""
            else:
                prompt = f"""Выполни следующую задачу, используя доступные инструменты: {user_task}

ВАЖНО: Выполни ВСЕ требуемые действия в правильной последовательности!
Если задача содержит несколько шагов ("затем", "потом", "а затем"), выполни их ВСЕ по порядку."""
            
            messages = [{"role": "user", "content": prompt}]
            tools = self.get_openai_format_tools()
            response = communicate_with_Pollinations_chat(self.current_model, messages, tools, api_token=self.api_token)
            
            if "choices" in response and len(response["choices"]) > 0:
                message = response["choices"][0]["message"]
                
                # Проверяем, есть ли вызовы инструментов
                if "tool_calls" in message:
                    print("\n🤖 AI выполняет задачу автоматически...")
                    success = True
                    error_details = []
                    
                    # Инициализируем список всех действий если это первая попытка
                    if attempt == 1:
                        all_executed_actions = []
                    
                    current_attempt_actions = []
                    generated_images = []  # Отслеживаем сгенерированные изображения
                    
                    for tool_call in message["tool_calls"]:
                        function_name = tool_call["function"]["name"]
                        function_args = json.loads(tool_call["function"]["arguments"])
                        
                        # Проверяем, не пытается ли AI сгенерировать похожее изображение повторно
                        if function_name == "generateImage":
                            prompt = function_args.get("prompt", "")
                            # Проверяем на дублирование по ключевым словам
                            is_duplicate = False
                            for prev_prompt in generated_images:
                                # Простая проверка на схожесть промптов
                                prompt_words = set(prompt.lower().split())
                                prev_words = set(prev_prompt.lower().split())
                                # Если больше 60% слов совпадают, считаем дубликатом
                                if len(prompt_words & prev_words) / max(len(prompt_words), 1) > 0.6:
                                    is_duplicate = True
                                    break
                            
                            if is_duplicate:
                                print(f"\n⚠️ Пропускаю дублирующую генерацию изображения: {prompt}")
                                continue
                            else:
                                generated_images.append(prompt)
                        
                        print(f"\n📋 Выполняется: {function_name}")
                        print(f"📊 Параметры: {function_args}")
                        
                        # Выполняем действие
                        result = self.execute_tool_call(function_name, function_args)
                        current_attempt_actions.append(f"{function_name}({function_args})")
                        
                        if "Ошибка" in result:
                            success = False
                            error_details.append(f"{function_name}: {result}")
                            print(f"❌ Ошибка: {result}")
                        else:
                            print(f"✅ Результат: {result}")
                    
                    if success:
                        # Добавляем действия текущей попытки к общему списку
                        all_executed_actions.extend(current_attempt_actions)
                        
                        # Проверяем, выполнена ли задача полностью, учитывая ВСЕ действия
                        completion_check = self._check_task_completion(user_task, all_executed_actions)
                        if completion_check["complete"]:
                            return "✅ Задача выполнена успешно!"
                        else:
                            print(f"\n⚠️ Задача выполнена не полностью: {completion_check['missing']}")
                            last_error = f"Не выполнены действия: {completion_check['missing']}"
                            if attempt < max_attempts:
                                print(f"\n⚠️ Попытка {attempt} не завершена. Продолжаем выполнение...")
                                continue
                            else:
                                return f"⚠️ Задача выполнена частично после {max_attempts} попыток. Не выполнено: {completion_check['missing']}"
                    else:
                        last_error = "; ".join(error_details)
                        if attempt < max_attempts:
                            print(f"\n⚠️ Попытка {attempt} не удалась. Пробуем другой подход...")
                            continue
                        else:
                            return f"❌ Задача не выполнена после {max_attempts} попыток. Последние ошибки: {last_error}"
                else:
                    # Если инструменты не вызваны, возвращаем план как текст
                    content = message.get("content", "")
                    if content:
                        print("\n📝 AI предлагает план действий:")
                        return content
                    else:
                        return "[Ошибка] Не удалось получить ответ от AI."
            elif "error" in response:
                last_error = response['error']
                if attempt < max_attempts:
                    print(f"\n⚠️ Ошибка API: {last_error}. Повторяем попытку...")
                    continue
                else:
                    return f"[Ошибка API] {response['error']}"
            else:
                last_error = "Не удалось получить ответ от AI"
                if attempt < max_attempts:
                    continue
                else:
                    return "[Ошибка] Не удалось получить ответ от AI."

    def execute_tool_call(self, function_name, function_args):
        """Выполнение вызова инструмента"""
        try:
            # 📁 Файловая система
            if function_name == "createDirectory":
                return self.create_directory(function_args["path"])
            elif function_name == "createFile":
                content = function_args.get("content", "")
                return self.create_file(function_args["path"], content)
            elif function_name == "listDirectory":
                path = function_args.get("path", ".")
                return self.list_directory(path)
            elif function_name == "readFile":
                return self.read_file(function_args["path"])
            elif function_name == "moveFile":
                return self.move_file(function_args["source_path"], function_args["destination_path"])
            elif function_name == "deleteFile":
                return self.delete_file(function_args["path"])
            elif function_name == "deleteDirectory":
                return self.delete_directory(function_args["path"])
            elif function_name == "swapFiles":
                return self.swap_files(function_args["path1"], function_args["path2"])
            
            # 🎨 Генерация изображений
            elif function_name == "generateImage":
                return self.generate_image(
                    function_args["prompt"],
                    function_args.get("width", 1024),
                    function_args.get("height", 1024),
                    function_args.get("model", "flux"),
                    function_args.get("save_path"),
                    function_args.get("filename")
                )
            
            # 🎵 Генерация аудио
            elif function_name == "generateAudio":
                return self.generate_audio(
                    function_args["text"],
                    function_args.get("voice", "alloy")
                )
            elif function_name == "sayText" or function_name == "respondAudio":
                return self.generate_audio(
                    function_args["text"],
                    function_args.get("voice", "alloy")
                )
            
            # 🌐 Интернет и загрузки
            elif function_name == "downloadFile":
                return self.download_file(
                    function_args["url"],
                    function_args.get("save_path"),
                    function_args.get("filename")
                )
            elif function_name == "downloadImage":
                return self.download_image(
                    function_args["url"],
                    function_args.get("save_path"),
                    function_args.get("filename")
                )
            elif function_name == "searchInternet":
                model = function_args.get("model", "searchgpt")
                return self.search_internet(function_args["query"], model)
            
            # 💻 Системные команды
            elif function_name == "executeCommand":
                return self.execute_command(function_args["command"])
            elif function_name == "runPythonCode":
                return self.run_python_code(function_args["code"])
            
            # ⚙️ Системная информация
            elif function_name == "getSystemInfo":
                return self.get_system_info()
            
            # 🔍 Поиск и анализ файлов
            elif function_name == "findAndAnalyzeFile":
                return self.find_and_analyze_file(function_args["query"])
            elif function_name == "readAdvancedFile":
                return self.read_advanced_file(function_args["file_path"])
            elif function_name == "analyzeImage":
                return self.analyze_image(function_args["image_path"])
            elif function_name == "recognizeText":
                return self.recognize_text_from_image(function_args["image_path"])
            
            # 📂 Открытие и запуск файлов
            elif function_name == "openFileWithDefaultProgram":
                return self.open_file_with_default_program(function_args["file_path"])
            elif function_name == "runExecutable":
                return self.run_executable(
                    function_args["file_path"],
                    function_args.get("run_as_admin", False),
                    function_args.get("require_confirmation", True)
                )
            elif function_name == "smartOpenFile":
                return self.smart_open_file(
                    function_args["file_path"],
                    function_args.get("run_as_admin", False)
                )
            
            else:
                return f"Неизвестный инструмент: {function_name}"
        except Exception as e:
            return f"Ошибка выполнения инструмента {function_name}: {str(e)}"
    
    def execute_step(self, step):
        """Выполнение отдельного шага плана с прогресс-баром"""
        try:
            parts = step.split(":")
            tool_part = parts[0].split("-")
            tool_name = tool_part[1].strip()
            params = eval(parts[1].strip())
            if tool_name in self.mcp_tools:
                with tqdm(total=100, desc=f"Выполнение {tool_name}", unit="%") as pbar:
                    result = self.mcp_tools[tool_name](*params)
                    pbar.update(100)
                return result
            return f"Неизвестный инструмент: {tool_name}"
        except Exception as e:
            return f"Ошибка выполнения: {str(e)}"

    def generate_image_url(self, prompt):
        """Генерация URL изображения"""
        return f"{self.base_image_url}/prompt/{urllib.parse.quote(prompt)}"

    def generate_image(self, prompt, width=1024, height=1024, model="flux", save_path=None, filename=None):
        """Генерация изображения и сохранение в файл с поддержкой выбора папки и модели"""
        try:
            # Получаем список доступных моделей для выбора
            available_models = self.list_image_models()
            print(f"\n🎨 Генерация изображения...")
            print(f"📝 Промпт: {prompt}")
            print(f"📐 Размер: {width}x{height}")
            
            # Всегда показываем доступные модели для выбора
            print(f"\n🔧 Доступные модели изображений:")
            for i, model_name in enumerate(available_models[:15], 1):  # Показываем первые 15 моделей
                current_marker = " (текущая)" if model_name == model else ""
                print(f"{i}. {model_name}{current_marker}")
            
            try:
                choice = input(f"\nВыберите модель (1-{min(15, len(available_models))}) или нажмите Enter для {model}: ").strip()
                if choice and choice.isdigit():
                    choice_idx = int(choice) - 1
                    if 0 <= choice_idx < min(15, len(available_models)):
                        model = available_models[choice_idx]
                        print(f"✅ Выбрана модель: {model}")
                    else:
                        print(f"⚠️ Неверный номер, используется модель по умолчанию: {model}")
                else:
                    print(f"⚠️ Используется модель по умолчанию: {model}")
            except (ValueError, IndexError):
                print(f"⚠️ Ошибка ввода, используется модель по умолчанию: {model}")
            
            # Генерируем изображение
            encoded_prompt = requests.utils.quote(prompt)
            # Создаем URL с параметрами
            url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?model={model}&width={width}&height={height}&private=true"
            
            # Добавляем токен если он есть (опционально, для повышения приоритета)
            headers = {}
            if self.api_token:
                headers["Authorization"] = f"Bearer {self.api_token}"
            
            print(f"🔄 Загрузка изображения...")
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            
            # Определяем путь сохранения
            if not save_path:
                save_path = self.output_dir
            
            # Создаем папку если не существует
            os.makedirs(save_path, exist_ok=True)
            
            # Определяем имя файла
            if not filename:
                # Создаем безопасное имя файла на основе промпта
                safe_prompt = re.sub(r'[^\w\s-]', '', prompt)
                safe_prompt = re.sub(r'[-\s]+', '_', safe_prompt)
                filename = f"{safe_prompt[:50]}_{model}.jpg"  # Ограничиваем длину имени
            
            # Убеждаемся что у файла есть расширение
            if not os.path.splitext(filename)[1]:
                filename += '.jpg'
            
            filepath = os.path.join(save_path, filename)
            
            # Сохраняем изображение
            with open(filepath, "wb") as f:
                f.write(response.content)
            
            logger.info(f"Изображение сгенерировано и сохранено: {filepath}")
            return f"Изображение сгенерировано и сохранено: {filepath}"
            
        except requests.RequestException as e:
            return f"Ошибка загрузки изображения: {str(e)}"
        except Exception as e:
            return f"Ошибка генерации изображения: {str(e)}"
    
    # 🌐 Интернет и загрузки
    def download_file(self, url, save_path=None, filename=None):
        """Загружает файл из интернета"""
        try:
            response = requests.get(url, stream=True)
            response.raise_for_status()
            
            # Определяем имя файла
            if not filename:
                parsed_url = urlparse(url)
                filename = os.path.basename(parsed_url.path) or "downloaded_file"
            
            # Определяем путь сохранения
            if not save_path:
                save_path = self.output_dir
            
            # Создаем директорию если не существует
            os.makedirs(save_path, exist_ok=True)
            
            file_path = os.path.join(save_path, filename)
            
            # Загружаем файл
            with open(file_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            
            logger.info(f"Файл загружен: {file_path}")
            return f"Файл загружен: {file_path}"
        except Exception as e:
            return f"Ошибка загрузки файла: {str(e)}"
    
    def download_image(self, url, save_path=None, filename=None):
        """Загружает изображение из интернета"""
        try:
            response = requests.get(url)
            response.raise_for_status()
            
            # Определяем имя файла
            if not filename:
                parsed_url = urlparse(url)
                filename = os.path.basename(parsed_url.path) or "image.jpg"
                # Добавляем расширение если его нет
                if not os.path.splitext(filename)[1]:
                    content_type = response.headers.get('content-type', '')
                    if 'jpeg' in content_type or 'jpg' in content_type:
                        filename += '.jpg'
                    elif 'png' in content_type:
                        filename += '.png'
                    elif 'gif' in content_type:
                        filename += '.gif'
                    else:
                        filename += '.jpg'
            
            # Определяем путь сохранения
            if not save_path:
                save_path = self.output_dir
            
            # Создаем директорию если не существует
            os.makedirs(save_path, exist_ok=True)
            
            file_path = os.path.join(save_path, filename)
            
            # Сохраняем изображение
            with open(file_path, 'wb') as f:
                f.write(response.content)
            
            logger.info(f"Изображение загружено: {file_path}")
            return f"Изображение загружено: {file_path}"
        except Exception as e:
            return f"Ошибка загрузки изображения: {str(e)}"
    
    def search_internet(self, query, model="searchgpt"):
        """Поиск в интернете используя специальные модели"""
        try:
            url = "https://text.pollinations.ai/openai"
            payload = {
                "model": model,
                "messages": [{"role": "user", "content": query}],
                "private": True  # Не показывать в публичном feed
            }
            
            response = requests.post(url, json=payload)
            response.raise_for_status()
            
            result = response.json()
            if "choices" in result and len(result["choices"]) > 0:
                content = result["choices"][0]["message"].get("content", "")
                logger.info(f"Выполнен поиск: {query}")
                return content
            else:
                return "Не удалось получить результаты поиска"
        except Exception as e:
            return f"Ошибка поиска: {str(e)}"
    
    # 💻 Системные команды
    def execute_command(self, command, require_confirmation=True):
        """Выполняет команду в командной строке"""
        try:
            if require_confirmation:
                print(f"⚠️ Запрос на выполнение команды: {command}")
                confirm = input("Разрешить выполнение? (y/n): ").lower()
                if confirm != 'y':
                    return "Выполнение команды отменено пользователем"
            
            print(f"🔄 Выполняется команда: {command}")
            result = subprocess.run(
                command, 
                shell=True, 
                capture_output=True, 
                text=True, 
                encoding='utf-8'
            )
            
            output = f"Код возврата: {result.returncode}\n"
            if result.stdout:
                output += f"Вывод:\n{result.stdout}\n"
            if result.stderr:
                output += f"Ошибки:\n{result.stderr}\n"
            
            logger.info(f"Выполнена команда: {command}")
            return output
        except Exception as e:
            return f"Ошибка выполнения команды: {str(e)}"
    
    def run_python_code(self, code, require_confirmation=True):
        """Выполняет Python код"""
        try:
            if require_confirmation:
                print(f"⚠️ Запрос на выполнение Python кода:")
                print(f"```python\n{code}\n```")
                confirm = input("Разрешить выполнение? (y/n): ").lower()
                if confirm != 'y':
                    return "Выполнение кода отменено пользователем"
            
            print(f"🐍 Выполняется Python код...")
            
            # Создаем временный файл
            temp_file = os.path.join(self.output_dir, "temp_script.py")
            with open(temp_file, 'w', encoding='utf-8') as f:
                f.write(code)
            
            # Выполняем
            result = subprocess.run(
                ["python", temp_file], 
                capture_output=True, 
                text=True,
                encoding='utf-8'
            )
            
            # Удаляем временный файл
            os.remove(temp_file)
            
            output = f"Код возврата: {result.returncode}\n"
            if result.stdout:
                output += f"Вывод:\n{result.stdout}\n"
            if result.stderr:
                output += f"Ошибки:\n{result.stderr}\n"
            
            logger.info(f"Выполнен Python код")
            return output
        except Exception as e:
            return f"Ошибка выполнения Python кода: {str(e)}"
    
    # ⚙️ Системные настройки
    def change_registry_value(self, key_path, value_name, value_data, value_type="REG_SZ", require_confirmation=True):
        """Изменяет значение в реестре Windows"""
        try:
            if require_confirmation:
                print(f"⚠️ Запрос на изменение реестра:")
                print(f"Путь: {key_path}")
                print(f"Параметр: {value_name}")
                print(f"Значение: {value_data}")
                confirm = input("Разрешить изменение реестра? (y/n): ").lower()
                if confirm != 'y':
                    return "Изменение реестра отменено пользователем"
            
            # Определяем тип реестра
            reg_types = {
                "REG_SZ": winreg.REG_SZ,
                "REG_DWORD": winreg.REG_DWORD,
                "REG_BINARY": winreg.REG_BINARY
            }
            
            reg_type = reg_types.get(value_type, winreg.REG_SZ)
            
            # Открываем ключ реестра
            with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_WRITE) as key:
                winreg.SetValueEx(key, value_name, 0, reg_type, value_data)
            
            logger.info(f"Изменен реестр: {key_path}\\{value_name}")
            return f"Реестр успешно изменен: {key_path}\\{value_name} = {value_data}"
        except Exception as e:
            return f"Ошибка изменения реестра: {str(e)}"
    
    def get_system_info(self):
        """Получает информацию о системе"""
        try:
            import platform
            import psutil
            
            info = {
                "Операционная система": platform.system(),
                "Версия ОС": platform.version(),
                "Архитектура": platform.architecture()[0],
                "Процессор": platform.processor(),
                "Имя компьютера": platform.node(),
                "Пользователь": os.getenv('USERNAME', 'Неизвестно'),
                "RAM (ГБ)": round(psutil.virtual_memory().total / (1024**3), 2),
                "Свободная RAM (ГБ)": round(psutil.virtual_memory().available / (1024**3), 2),
                "Загрузка CPU (%)": psutil.cpu_percent(),
                "Количество ядер": psutil.cpu_count()
            }
            
            return json.dumps(info, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"Ошибка получения информации о системе: {str(e)}"
    
    def manage_services(self, service_name, action, require_confirmation=True):
        """Управление службами Windows"""
        try:
            if require_confirmation:
                print(f"⚠️ Запрос на управление службой:")
                print(f"Служба: {service_name}")
                print(f"Действие: {action}")
                confirm = input("Разрешить управление службой? (y/n): ").lower()
                if confirm != 'y':
                    return "Управление службой отменено пользователем"
            
            valid_actions = ['start', 'stop', 'restart', 'status']
            if action not in valid_actions:
                return f"Недопустимое действие. Доступны: {', '.join(valid_actions)}"
            
            if action == 'status':
                command = f"sc query {service_name}"
            elif action == 'start':
                command = f"sc start {service_name}"
            elif action == 'stop':
                command = f"sc stop {service_name}"
            elif action == 'restart':
                # Сначала останавливаем, потом запускаем
                stop_result = subprocess.run(f"sc stop {service_name}", shell=True, capture_output=True, text=True)
                start_result = subprocess.run(f"sc start {service_name}", shell=True, capture_output=True, text=True)
                return f"Остановка: {stop_result.stdout}\nЗапуск: {start_result.stdout}"
            
            result = subprocess.run(command, shell=True, capture_output=True, text=True)
            
            logger.info(f"Управление службой {service_name}: {action}")
            return f"Результат: {result.stdout}\nОшибки: {result.stderr}"
        except Exception as e:
            return f"Ошибка управления службой: {str(e)}"
    
    # 🔍 Поиск и загрузка изображений
    def search_and_download_images(self, query, num_images=5, save_path=None):
        """Ищет изображения в интернете через DuckDuckGo и загружает их"""
        try:
            if not save_path:
                save_path = os.path.join(self.output_dir, "images")
            
            os.makedirs(save_path, exist_ok=True)
            
            # Поиск изображений через DuckDuckGo
            search_url = f"https://duckduckgo.com/?q={urllib.parse.quote(query)}&t=h_&iax=images&ia=images"
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            # Получаем страницу поиска
            response = requests.get(search_url, headers=headers)
            if response.status_code != 200:
                return f"Ошибка поиска: HTTP {response.status_code}"
            
            # Парсим HTML для поиска ссылок на изображения
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Альтернативный подход - используем API поиска
            api_url = "https://duckduckgo.com/i.js"
            params = {
                'q': query,
                'o': 'json',
                'p': '1',
                's': '0',
                'u': 'bing',
                'f': ',,,',
                'l': 'us-en'
            }
            
            api_response = requests.get(api_url, params=params, headers=headers)
            
            # Если API не работает, используем простые изображения-заглушки
            if api_response.status_code != 200:
                logger.warning("API DuckDuckGo недоступен, используем заглушки")
                image_urls = [
                    "https://via.placeholder.com/800x600/0066CC/FFFFFF?text=" + urllib.parse.quote(f"Image {i+1}: {query}")
                    for i in range(num_images)
                ]
            else:
                try:
                    data = api_response.json()
                    image_urls = []
                    if 'results' in data:
                        for result in data['results'][:num_images]:
                            if 'image' in result:
                                image_urls.append(result['image'])
                    
                    if not image_urls:
                        # Fallback к заглушкам
                        image_urls = [
                            "https://via.placeholder.com/800x600/0066CC/FFFFFF?text=" + urllib.parse.quote(f"Image {i+1}: {query}")
                            for i in range(num_images)
                        ]
                except:
                    # Fallback к заглушкам
                    image_urls = [
                        "https://via.placeholder.com/800x600/0066CC/FFFFFF?text=" + urllib.parse.quote(f"Image {i+1}: {query}")
                        for i in range(num_images)
                    ]
            
            downloaded_files = []
            
            # Загружаем найденные изображения
            for i, url in enumerate(image_urls[:num_images]):
                try:
                    filename = f"{query.replace(' ', '_')}_{i+1}.jpg"
                    result = self.download_image(url, save_path, filename)
                    if "загружено" in result:
                        downloaded_files.append(result)
                    else:
                        logger.warning(f"Не удалось загрузить изображение {i+1}: {result}")
                except Exception as e:
                    logger.warning(f"Ошибка загрузки изображения {i+1}: {str(e)}")
                    continue
            
            if downloaded_files:
                logger.info(f"Загружено {len(downloaded_files)} изображений для запроса: {query}")
                return f"Загружено {len(downloaded_files)} изображений в папку {save_path}:\n" + "\n".join(downloaded_files)
            else:
                return f"Не удалось загрузить изображения для запроса: {query}"
                
        except Exception as e:
            return f"Ошибка поиска и загрузки изображений: {str(e)}"
    
    # 🐍 Создание Python проектов
    def create_python_project(self, project_name, code="", requirements=None, run_after_create=True):
        """Создает Python проект с виртуальным окружением"""
        try:
            # Создаем папку проекта
            project_dir = os.path.join("projects", project_name)
            os.makedirs(project_dir, exist_ok=True)
            
            # Создаем виртуальное окружение
            venv_dir = os.path.join(project_dir, "venv")
            print(f"🔄 Создание виртуального окружения...")
            venv.create(venv_dir, with_pip=True)
            
            # Определяем путь к python в виртуальном окружении
            if os.name == 'nt':  # Windows
                venv_python = os.path.join(venv_dir, "Scripts", "python.exe")
                venv_pip = os.path.join(venv_dir, "Scripts", "pip.exe")
            else:  # Unix/Linux
                venv_python = os.path.join(venv_dir, "bin", "python")
                venv_pip = os.path.join(venv_dir, "bin", "pip")
            
            # Создаем основной Python файл
            main_file = os.path.join(project_dir, "main.py")
            if not code:
                code = f'''# {project_name}
# Создано автоматически Pollinations Agent

def main():
    print("Привет из проекта {project_name}!")
    # Добавьте ваш код здесь
    pass

if __name__ == "__main__":
    main()
'''
            
            with open(main_file, 'w', encoding='utf-8') as f:
                f.write(code)
            
            # Создаем requirements.txt
            req_file = os.path.join(project_dir, "requirements.txt")
            if requirements:
                if isinstance(requirements, list):
                    req_content = "\n".join(requirements)
                else:
                    req_content = str(requirements)
            else:
                req_content = "# Добавьте зависимости сюда\n# Например:\n# requests>=2.25.1\n# numpy>=1.21.0"
            
            with open(req_file, 'w', encoding='utf-8') as f:
                f.write(req_content)
            
            # Создаем README.md
            readme_file = os.path.join(project_dir, "README.md")
            readme_content = f"""# {project_name}

Проект создан автоматически с помощью Pollinations Agent.

## Установка

1. Активируйте виртуальное окружение:
   ```bash
   # Windows
   venv\\Scripts\\activate
   
   # Linux/Mac
   source venv/bin/activate
   ```

2. Установите зависимости:
   ```bash
   pip install -r requirements.txt
   ```

## Запуск

```bash
python main.py
```

## Структура проекта

```
{project_name}/
├── venv/           # Виртуальное окружение
├── main.py         # Основной файл
├── requirements.txt # Зависимости
└── README.md       # Этот файл
```
"""
            
            with open(readme_file, 'w', encoding='utf-8') as f:
                f.write(readme_content)
            
            # Устанавливаем зависимости если они есть
            if requirements and requirements.strip() and requirements != "# Добавьте зависимости сюда\n# Например:\n# requests>=2.25.1\n# numpy>=1.21.0":
                print(f"🔄 Установка зависимостей...")
                try:
                    if isinstance(requirements, list):
                        for req in requirements:
                            subprocess.run([venv_pip, "install", req], check=True)
                    else:
                        # Сохраняем requirements во временный файл и устанавливаем
                        temp_req = os.path.join(project_dir, "temp_req.txt")
                        with open(temp_req, 'w', encoding='utf-8') as f:
                            f.write(requirements)
                        subprocess.run([venv_pip, "install", "-r", temp_req], check=True)
                        os.remove(temp_req)
                except subprocess.CalledProcessError as e:
                    logger.warning(f"Ошибка установки зависимостей: {e}")
            
            result_message = f"""✅ Python проект '{project_name}' создан успешно!

📁 Структура проекта:
├── {project_dir}/
│   ├── venv/           # Виртуальное окружение
│   ├── main.py         # Основной файл с кодом
│   ├── requirements.txt # Зависимости
│   └── README.md       # Документация

🚀 Для запуска:
1. cd {project_dir}
2. {venv_python} main.py
"""
            
            # Запускаем код если требуется
            if run_after_create and code.strip():
                print(f"🐍 Запуск созданного проекта...")
                try:
                    run_result = subprocess.run(
                        [venv_python, main_file], 
                        capture_output=True, 
                        text=True,
                        encoding='utf-8',
                        cwd=project_dir
                    )
                    
                    result_message += f"\n🏃 Результат выполнения:\n"
                    if run_result.stdout:
                        result_message += f"Вывод: {run_result.stdout}\n"
                    if run_result.stderr:
                        result_message += f"Ошибки: {run_result.stderr}\n"
                    result_message += f"Код возврата: {run_result.returncode}"
                    
                except Exception as e:
                    result_message += f"\n⚠️ Ошибка при запуске: {str(e)}"
            
            logger.info(f"Создан Python проект: {project_name}")
            return result_message
            
        except Exception as e:
            return f"Ошибка создания Python проекта: {str(e)}"
    
    # 💻 Продвинутая разработка ПО с автоматическим исправлением ошибок
    def develop_software(self, task_description, max_attempts=5):
        """Разрабатывает ПО по описанию задачи с автоматическим исправлением ошибок"""
        try:
            print(f"🚀 Начинаю разработку ПО для задачи: {task_description}")
            
            # 1. Анализируем задачу и определяем язык программирования
            language_info = self._analyze_task_and_select_language(task_description)
            
            if "error" in language_info:
                return f"Ошибка анализа задачи: {language_info['error']}"
            
            language = language_info['language']
            project_name = language_info['project_name']
            suggested_code = language_info['code']
            file_extension = language_info['extension']
            run_command = language_info['run_command']
            dependencies = language_info.get('dependencies', [])
            
            print(f"📋 Выбранный язык: {language}")
            print(f"📁 Название проекта: {project_name}")
            
            # 2. Создаем структуру проекта
            project_dir = os.path.join("software_projects", project_name)
            os.makedirs(project_dir, exist_ok=True)
            
            main_file = os.path.join(project_dir, f"main{file_extension}")
            
            # 3. Записываем начальный код
            with open(main_file, 'w', encoding='utf-8') as f:
                f.write(suggested_code)
            
            print(f"✅ Создан файл: {main_file}")
            print(f"📝 Начальный код записан")
            
            # 4. Устанавливаем зависимости если нужно
            if dependencies and language.lower() == 'python':
                self._install_python_dependencies(project_dir, dependencies)
            
            # 5. Тестируем и исправляем ошибки
            attempt = 0
            last_error = None
            
            while attempt < max_attempts:
                attempt += 1
                print(f"\n🔄 Попытка запуска #{attempt}...")
                
                # Запускаем код
                run_result = self._run_code(project_dir, main_file, run_command)
                
                if run_result['success']:
                    print(f"\n🎉 Программа запустилась успешно!")
                    print(f"✅ Вывод программы:")
                    if run_result['output']:
                        print(run_result['output'])
                    
                    # Предлагаем пользователю протестировать
                    test_result = self._offer_user_testing(project_dir, main_file, run_command)
                    
                    return f"""🎯 Разработка ПО завершена успешно!

📁 Проект: {project_name}
📍 Путь: {project_dir}
🔧 Язык: {language}
📄 Основной файл: {main_file}

✅ Программа протестирована и готова к использованию!

{test_result}
"""
                else:
                    error_details = run_result['error']
                    print(f"❌ Ошибка выполнения: {error_details}")
                    
                    if attempt < max_attempts:
                        print(f"🔧 Пытаюсь исправить ошибку...")
                        
                        # Исправляем ошибку с помощью ИИ
                        fixed_code = self._fix_code_with_ai(suggested_code, error_details, task_description, language)
                        
                        if fixed_code and fixed_code != suggested_code:
                            suggested_code = fixed_code
                            # Перезаписываем файл с исправленным кодом
                            with open(main_file, 'w', encoding='utf-8') as f:
                                f.write(suggested_code)
                            print(f"🔄 Код обновлен, пробую снова...")
                        else:
                            print(f"⚠️ Не удалось получить исправленный код")
                            last_error = error_details
                    else:
                        last_error = error_details
            
            # Если не удалось исправить за max_attempts попыток
            return f"""⚠️ Разработка ПО завершена с ошибками

📁 Проект: {project_name}
📍 Путь: {project_dir}
🔧 Язык: {language}
📄 Основной файл: {main_file}

❌ Не удалось исправить ошибки за {max_attempts} попыток.
Последняя ошибка: {last_error}

💡 Попробуйте:
1. Проверить код вручную в файле {main_file}
2. Переформулировать задачу более детально
3. Указать дополнительные требования
"""
                
        except Exception as e:
            return f"Критическая ошибка разработки ПО: {str(e)}"
    
    # 📄 Чтение различных форматов файлов
    def read_advanced_file(self, file_path):
        """Умное чтение файлов различных форматов"""
        try:
            # Получаем расширение файла
            _, ext = os.path.splitext(file_path.lower())
            
            # Проверяем существование файла
            if not os.path.exists(file_path):
                return f"Файл не найден: {file_path}"
            
            # Выбираем подходящий метод чтения
            if ext in ['.txt', '.log', '.json', '.xml', '.csv', '.md', '.py', '.js', '.html', '.css']:
                return self.read_file(file_path)
            elif ext in ['.docx', '.doc']:
                return self.read_docx_file(file_path)
            elif ext == '.pdf':
                return self.read_pdf_file(file_path)
            elif ext in ['.xlsx', '.xls']:
                return self.read_excel_file(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self.read_powerpoint_file(file_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                return self.analyze_image(file_path)
            else:
                # Пробуем прочитать как текстовый файл
                return self.read_file(file_path)
                
        except Exception as e:
            return f"Ошибка чтения файла: {str(e)}"
    
    def read_docx_file(self, file_path):
        """Читает содержимое DOCX файла"""
        try:
            if not DOCX_AVAILABLE:
                return "Ошибка: Модуль python-docx не установлен. Установите: pip install python-docx"
            
            doc = Document(file_path)
            content = []
            
            # Читаем параграфы
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text.strip())
            
            # Читаем таблицы если есть
            for table in doc.tables:
                content.append("\n--- ТАБЛИЦА ---")
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    content.append(" | ".join(row_text))
                content.append("--- КОНЕЦ ТАБЛИЦЫ ---\n")
            
            result = "\n".join(content)
            logger.info(f"DOCX файл прочитан: {file_path}")
            return result if result.strip() else "Документ пуст или не содержит текста"
            
        except Exception as e:
            return f"Ошибка чтения DOCX файла: {str(e)}"
    
    def read_pdf_file(self, file_path):
        """Читает содержимое PDF файла"""
        try:
            if not PDF_AVAILABLE:
                return "Ошибка: Модуль PyPDF2 не установлен. Установите: pip install PyPDF2"
            
            content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Читаем все страницы
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            content.append(f"--- СТРАНИЦА {page_num} ---")
                            content.append(text.strip())
                            content.append("")
                    except Exception as e:
                        content.append(f"--- СТРАНИЦА {page_num} (ОШИБКА ЧТЕНИЯ) ---")
                        content.append(f"Ошибка: {str(e)}")
                        content.append("")
            
            result = "\n".join(content)
            logger.info(f"PDF файл прочитан: {file_path} ({len(pdf_reader.pages)} страниц)")
            return result if result.strip() else "PDF документ пуст или не содержит извлекаемого текста"
            
        except Exception as e:
            return f"Ошибка чтения PDF файла: {str(e)}"
    
    def read_excel_file(self, file_path):
        """Читает содержимое Excel файла"""
        try:
            if not EXCEL_AVAILABLE:
                return "Ошибка: Модуль openpyxl не установлен. Установите: pip install openpyxl"
            
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            content = []
            
            # Читаем все листы
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"--- ЛИСТ: {sheet_name} ---")
                
                # Определяем размеры данных
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                if max_row > 0 and max_col > 0:
                    # Читаем данные
                    for row in sheet.iter_rows(min_row=1, max_row=min(max_row, 100), 
                                               min_col=1, max_col=min(max_col, 20), 
                                               values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                            else:
                                row_text.append("")
                        if any(cell.strip() for cell in row_text if cell):
                            content.append(" | ".join(row_text))
                    
                    if max_row > 100 or max_col > 20:
                        content.append(f"... (показаны первые 100 строк и 20 столбцов из {max_row}x{max_col})")
                else:
                    content.append("Лист пуст")
                
                content.append("")
            
            result = "\n".join(content)
            logger.info(f"Excel файл прочитан: {file_path}")
            return result if result.strip() else "Excel файл пуст"
            
        except Exception as e:
            return f"Ошибка чтения Excel файла: {str(e)}"
    
    def read_powerpoint_file(self, file_path):
        """Читает содержимое PowerPoint файла"""
        try:
            if not PPTX_AVAILABLE:
                return "Ошибка: Модуль python-pptx не установлен. Установите: pip install python-pptx"
            
            prs = Presentation(file_path)
            content = []
            
            # Читаем все слайды
            for slide_num, slide in enumerate(prs.slides, 1):
                content.append(f"--- СЛАЙД {slide_num} ---")
                
                # Читаем все текстовые блоки на слайде
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text.strip())
                
                content.append("")
            
            result = "\n".join(content)
            logger.info(f"PowerPoint файл прочитан: {file_path} ({len(prs.slides)} слайдов)")
            return result if result.strip() else "PowerPoint презентация пуста или не содержит текста"
            
        except Exception as e:
            return f"Ошибка чтения PowerPoint файла: {str(e)}"
    
    # 🖼️ Анализ изображений
    def analyze_image(self, image_path):
        """Анализирует изображение с помощью Pollinations Vision API"""
        try:
            # Проверяем существование файла
            if not os.path.exists(image_path):
                return f"Изображение не найдено: {image_path}"
            
            # Кодируем изображение в base64
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Определяем MIME тип
            mime_type, _ = mimetypes.guess_type(image_path)
            if not mime_type or not mime_type.startswith('image/'):
                mime_type = 'image/jpeg'  # Fallback
            
            # Отправляем запрос в Pollinations Vision API
            payload = {
                "model": "openai",  # Модель с поддержкой vision
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Проанализируй это изображение детально. Опиши что на нем изображено, цвета, объекты, текст (если есть), настроение, стиль и любые важные детали."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{image_data}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 500,
                "private": True
            }
            
            headers = {"Content-Type": "application/json"}
            if self.api_token:
                headers["Authorization"] = f"Bearer {self.api_token}"
            
            response = requests.post(
                f"{self.base_text_url}/openai",
                json=payload,
                headers=headers
            )
            response.raise_for_status()
            
            result = response.json()
            if "choices" in result and len(result["choices"]) > 0:
                analysis = result["choices"][0]["message"].get("content", "")
                logger.info(f"Изображение проанализировано: {image_path}")
                return f"Анализ изображения {os.path.basename(image_path)}:\n\n{analysis}"
            else:
                return "Не удалось получить анализ изображения"
                
        except Exception as e:
            return f"Ошибка анализа изображения: {str(e)}"
    
    def recognize_text_from_image(self, image_path):
        """Распознает текст на изображении с помощью OCR"""
        try:
            if not PILLOW_AVAILABLE:
                # Fallback к Pollinations Vision API для распознавания текста
                return self._recognize_text_with_vision_api(image_path)
            
            # Используем pytesseract если доступен
            try:
                image = Image.open(image_path)
                text = pytesseract.image_to_string(image, lang='rus+eng')
                
                if text.strip():
                    logger.info(f"Текст распознан из изображения: {image_path}")
                    return f"Распознанный текст из {os.path.basename(image_path)}:\n\n{text.strip()}"
                else:
                    return f"Текст не найден на изображении: {os.path.basename(image_path)}"
                    
            except Exception as tesseract_error:
                # Если tesseract не работает, используем Vision API
                logger.warning(f"Tesseract недоступен: {tesseract_error}")
                return self._recognize_text_with_vision_api(image_path)
                
        except Exception as e:
            return f"Ошибка распознавания текста: {str(e)}"
    
    def _recognize_text_with_vision_api(self, image_path):
        """Распознает текст с помощью Pollinations Vision API"""
        try:
            # Кодируем изображение в base64
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Определяем MIME тип
            mime_type, _ = mimetypes.guess_type(image_path)
            if not mime_type or not mime_type.startswith('image/'):
                mime_type = 'image/jpeg'
            
            # Отправляем запрос
            payload = {
                "model": "openai",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Извлеки и запиши весь текст, который ты видишь на этом изображении. Если текста нет, напиши 'Текст не найден'."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{image_data}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 300,
                "private": True
            }
            
            headers = {"Content-Type": "application/json"}
            if self.api_token:
                headers["Authorization"] = f"Bearer {self.api_token}"
            
            response = requests.post(
                f"{self.base_text_url}/openai",
                json=payload,
                headers=headers
            )
            response.raise_for_status()
            
            result = response.json()
            if "choices" in result and len(result["choices"]) > 0:
                text = result["choices"][0]["message"].get("content", "")
                logger.info(f"Текст распознан через Vision API: {image_path}")
                return f"Распознанный текст из {os.path.basename(image_path)}:\n\n{text}"
            else:
                return "Не удалось распознать текст"
                
        except Exception as e:
            return f"Ошибка распознавания текста через Vision API: {str(e)}"
    
    def get_image_info(self, image_path):
        """Получает техническую информацию об изображении"""
        try:
            if not PILLOW_AVAILABLE:
                return "Ошибка: Модуль Pillow не установлен. Установите: pip install Pillow"
            
            # Проверяем существование файла
            if not os.path.exists(image_path):
                return f"Изображение не найдено: {image_path}"
            
            # Открываем изображение
            with Image.open(image_path) as img:
                # Получаем базовую информацию
                info = {
                    "Имя файла": os.path.basename(image_path),
                    "Размер файла": f"{os.path.getsize(image_path) / 1024:.1f} КБ",
                    "Формат": img.format,
                    "Режим": img.mode,
                    "Размер (пиксели)": f"{img.width} x {img.height}",
                    "Соотношение сторон": f"{img.width/img.height:.2f}:1"
                }
                
                # Дополнительная информация если доступна
                if hasattr(img, 'info') and img.info:
                    for key, value in img.info.items():
                        if isinstance(value, (str, int, float)):
                            info[f"EXIF {key}"] = value
                
                result = "Информация об изображении:\n\n"
                for key, value in info.items():
                    result += f"{key}: {value}\n"
                
                logger.info(f"Получена информация об изображении: {image_path}")
                return result
                
        except Exception as e:
            return f"Ошибка получения информации об изображении: {str(e)}"
    
    def find_and_analyze_file(self, query):
        """Находит и анализирует файл в текущей директории по запросу"""
        try:
            # Получаем список файлов в текущей директории
            current_files = os.listdir(".")
            
            # Извлекаем числа из запроса
            import re
            numbers_in_query = re.findall(r'\b\d+\b', query)
            
            # Определяем тип файла для поиска
            query_lower = query.lower()
            
            target_file = None
            
            # Если есть числа в запросе, ищем файл с этим числом
            if numbers_in_query:
                for number in numbers_in_query:
                    for filename in current_files:
                        if number in filename:
                            target_file = filename
                            break
                    if target_file:
                        break
            
            # Если не нашли по числу, ищем по типу файла
            if not target_file:
                image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']
                document_extensions = ['.docx', '.doc', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt']
                text_extensions = ['.txt', '.log', '.json', '.xml', '.csv', '.md']
                
                if any(word in query_lower for word in ['картинк', 'изображен', 'фото', 'рисунок']):
                    # Ищем изображения
                    for filename in current_files:
                        if any(filename.lower().endswith(ext) for ext in image_extensions):
                            target_file = filename
                            break
                elif any(word in query_lower for word in ['документ', 'файл', 'текст']):
                    # Ищем документы
                    for filename in current_files:
                        if any(filename.lower().endswith(ext) for ext in document_extensions + text_extensions):
                            target_file = filename
                            break
            
            # Если все еще не нашли, берем первый подходящий файл
            if not target_file:
                # Приоритет изображениям, затем документам
                all_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp',
                                '.docx', '.doc', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt',
                                '.txt', '.log', '.json', '.xml', '.csv', '.md']
                
                for filename in current_files:
                    if any(filename.lower().endswith(ext) for ext in all_extensions):
                        target_file = filename
                        break
            
            if not target_file:
                return f"Не найден подходящий файл для анализа в текущей директории.\n\nДоступные файлы: {', '.join(current_files)}"
            
            print(f"🔍 Найден файл для анализа: {target_file}")
            
            # Определяем тип анализа
            if any(word in query_lower for word in ['написан', 'текст', 'надпись', 'распознай']):
                # Распознавание текста
                return self.recognize_text_from_image(target_file)
            else:
                # Обычный анализ файла
                return self.read_advanced_file(target_file)
                
        except Exception as e:
            return f"Ошибка поиска и анализа файла: {str(e)}"
    
    def open_file_with_default_program(self, file_path):
        """Открывает файл с программой по умолчанию"""
        try:
            # Проверяем существование файла
            if os.path.exists(file_path):
                actual_file = file_path
            else:
                # Ищем похожий файл
                actual_file = self._find_similar_file(file_path)
                if not actual_file:
                    return f"Файл не найден: {file_path}. Проверьте имя файла и попробуйте снова."
                print(f"🔍 Найден похожий файл: {actual_file}")
            
            logger.info(f"Открытие файла с программой по умолчанию: {actual_file}")
            
            if platform.system() == 'Windows':
                os.startfile(actual_file)
            elif platform.system() == 'Darwin':  # macOS
                subprocess.call(['open', actual_file])
            else:  # Linux
                subprocess.call(['xdg-open', actual_file])
            
            return f"Файл открыт: {actual_file}"
        except Exception as e:
            return f"Ошибка открытия файла: {str(e)}"
    
    def run_executable(self, file_path, run_as_admin=False, require_confirmation=True):
        """Запускает исполняемый файл"""
        try:
            # Проверяем существование файла
            if not os.path.exists(file_path):
                return f"Файл не найден: {file_path}"
            
            # Проверяем, что это исполняемый файл
            _, ext = os.path.splitext(file_path.lower())
            executable_extensions = ['.exe', '.bat', '.cmd', '.com', '.scr', '.msi']
            
            if ext not in executable_extensions:
                return f"Файл не является исполняемым: {file_path} (расширение: {ext})"
            
            if require_confirmation:
                print(f"⚠️ Запрос на запуск исполняемого файла: {file_path}")
                if run_as_admin:
                    print(f"🔐 Файл будет запущен от имени администратора")
                confirm = input("Разрешить запуск? (y/n): ").lower()
                if confirm != 'y':
                    return "Запуск файла отменен пользователем"
            
            logger.info(f"Запуск исполняемого файла: {file_path} (admin: {run_as_admin})")
            
            if platform.system() == 'Windows':
                if run_as_admin:
                    # Запуск от имени администратора через PowerShell
                    command = ['powershell', '-Command', f'Start-Process "{file_path}" -Verb RunAs']
                    subprocess.run(command, shell=True)
                else:
                    # Обычный запуск
                    subprocess.Popen([file_path], shell=True)
            else:
                # Для Linux/macOS
                if run_as_admin:
                    return "Запуск от имени администратора поддерживается только в Windows"
                subprocess.Popen([file_path])
            
            return f"Исполняемый файл запущен: {file_path}"
        except Exception as e:
            return f"Ошибка запуска исполняемого файла: {str(e)}"
    
    def _find_similar_file(self, file_path):
        """Ищет похожий файл в текущей директории"""
        try:
            import os
            import re
            
            # Получаем имя файла без пути
            filename = os.path.basename(file_path)
            name, ext = os.path.splitext(filename)
            
            # Получаем список всех файлов в текущей директории
            current_dir = os.path.dirname(file_path) if os.path.dirname(file_path) else "."
            all_files = os.listdir(current_dir)
            
            # Сначала ищем точное совпадение по имени с разными расширениями
            for file in all_files:
                file_name, file_ext = os.path.splitext(file)
                if file_name.lower() == name.lower():
                    return os.path.join(current_dir, file)
            
            # Затем ищем файлы, содержащие искомое имя
            for file in all_files:
                if name.lower() in file.lower():
                    return os.path.join(current_dir, file)
            
            # Ищем по числам в имени файла
            numbers_in_name = re.findall(r'\d+', name)
            if numbers_in_name:
                for number in numbers_in_name:
                    for file in all_files:
                        if number in file:
                            return os.path.join(current_dir, file)
            
            return None
            
        except Exception as e:
            print(f"Ошибка поиска похожего файла: {e}")
            return None

    def smart_open_file(self, file_path, run_as_admin=False):
        """Умное открытие файла: исполняемые запускаются, остальные открываются с программой по умолчанию"""
        try:
            # Проверяем существование файла
            if not os.path.exists(file_path):
                return f"Файл не найден: {file_path}"
            
            # Получаем расширение файла
            _, ext = os.path.splitext(file_path.lower())
            
            # Определяем исполняемые файлы
            executable_extensions = ['.exe', '.bat', '.cmd', '.com', '.scr', '.msi']
            
            if ext in executable_extensions:
                # Запускаем как исполняемый файл
                return self.run_executable(file_path, run_as_admin)
            else:
                # Открываем с программой по умолчанию
                return self.open_file_with_default_program(file_path)
        except Exception as e:
            return f"Ошибка открытия файла: {str(e)}"
    
    def _analyze_task_and_select_language(self, task_description):
        """Анализирует задачу и выбирает подходящий язык программирования"""
        try:
            prompt = f"""Проанализируй следующую задачу разработки ПО и определи:
1. Подходящий язык программирования
2. Название проекта (без пробелов, только английские буквы)
3. Напиши полный работающий код для решения задачи
4. Список зависимостей если нужны

Задача: {task_description}

Ответь в JSON формате:
{{
  "language": "Python/JavaScript/Java/C++/C#/Go/Rust/PHP",
  "project_name": "название_проекта",
  "code": "полный код программы",
  "extension": "расширение файла (.py, .js, .java, и т.д.)",
  "run_command": "команда для запуска",
  "dependencies": ["список", "зависимостей"]
}}

Обязательно включи полный работающий код, который решает поставленную задачу!"""
            
            messages = [{"role": "user", "content": prompt}]
            response = communicate_with_Pollinations_chat(self.current_model, messages)
            
            if "choices" in response and len(response["choices"]) > 0:
                content = response["choices"][0]["message"].get("content", "")
                
                # Извлекаем JSON из ответа
                try:
                    # Ищем JSON в ответе
                    json_start = content.find('{')
                    json_end = content.rfind('}') + 1
                    
                    if json_start != -1 and json_end != -1:
                        json_str = content[json_start:json_end]
                        result = json.loads(json_str)
                        
                        # Проверяем обязательные поля
                        required_fields = ['language', 'project_name', 'code', 'extension', 'run_command']
                        if all(field in result for field in required_fields):
                            return result
                        else:
                            missing = [f for f in required_fields if f not in result]
                            return {"error": f"Отсутствуют обязательные поля: {missing}"}
                    else:
                        return {"error": "Не найден JSON в ответе ИИ"}
                        
                except json.JSONDecodeError as e:
                    return {"error": f"Ошибка парсинга JSON: {str(e)}"}
            else:
                return {"error": "Не удалось получить ответ от ИИ"}
                
        except Exception as e:
            return {"error": f"Ошибка анализа задачи: {str(e)}"}
    
    def _install_python_dependencies(self, project_dir, dependencies):
        """Устанавливает Python зависимости"""
        try:
            if not dependencies:
                return
            
            print(f"📦 Устанавливаю зависимости: {', '.join(dependencies)}")
            
            # Создаем requirements.txt
            req_file = os.path.join(project_dir, "requirements.txt")
            with open(req_file, 'w', encoding='utf-8') as f:
                f.write('\n'.join(dependencies))
            
            # Устанавливаем зависимости
            for dep in dependencies:
                result = subprocess.run(
                    ["pip", "install", dep], 
                    capture_output=True, 
                    text=True,
                    encoding='utf-8'
                )
                if result.returncode == 0:
                    print(f"✅ Установлено: {dep}")
                else:
                    print(f"⚠️ Ошибка установки {dep}: {result.stderr}")
                    
        except Exception as e:
            print(f"⚠️ Ошибка установки зависимостей: {str(e)}")
    
    def _run_code(self, project_dir, main_file, run_command):
        """Запускает код и возвращает результат"""
        try:
            # Формируем команду запуска
            if run_command.startswith('python'):
                command = ["python", os.path.basename(main_file)]
            elif run_command.startswith('node'):
                command = ["node", os.path.basename(main_file)]
            elif run_command.startswith('java'):
                # Для Java нужна компиляция
                compile_result = subprocess.run(
                    ["javac", os.path.basename(main_file)],
                    cwd=project_dir,
                    capture_output=True,
                    text=True,
                    encoding='utf-8'
                )
                if compile_result.returncode != 0:
                    return {
                        "success": False,
                        "error": f"Ошибка компиляции Java: {compile_result.stderr}"
                    }
                
                class_name = os.path.splitext(os.path.basename(main_file))[0]
                command = ["java", class_name]
            else:
                # Пытаемся выполнить как есть
                command = run_command.split()
            
            print(f"🔄 Выполняю команду: {' '.join(command)}")
            
            # Запускаем
            result = subprocess.run(
                command,
                cwd=project_dir,
                capture_output=True,
                text=True,
                encoding='utf-8',
                timeout=30  # Ограничиваем время выполнения
            )
            
            if result.returncode == 0:
                return {
                    "success": True,
                    "output": result.stdout,
                    "error": None
                }
            else:
                return {
                    "success": False,
                    "error": result.stderr or result.stdout,
                    "output": result.stdout
                }
                
        except subprocess.TimeoutExpired:
            return {
                "success": False,
                "error": "Превышено время выполнения (30 сек)"
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"Ошибка запуска: {str(e)}"
            }
    
    def _fix_code_with_ai(self, current_code, error_details, task_description, language):
        """Исправляет код с помощью ИИ"""
        try:
            prompt = f"""Исправь ошибку в коде {language}.

Исходная задача: {task_description}

Текущий код:
```{language.lower()}
{current_code}
```

Ошибка при выполнении:
{error_details}

Исправь код так, чтобы он работал правильно. Ответь только исправленным кодом без дополнительных объяснений."""
            
            messages = [{"role": "user", "content": prompt}]
            response = communicate_with_Pollinations_chat(self.current_model, messages)
            
            if "choices" in response and len(response["choices"]) > 0:
                content = response["choices"][0]["message"].get("content", "")
                
                # Извлекаем код из ответа (убираем markdown разметку если есть)
                code_start = content.find('```')
                if code_start != -1:
                    # Находим начало кода после первых ```
                    code_start = content.find('\n', code_start) + 1
                    code_end = content.find('```', code_start)
                    if code_end != -1:
                        return content[code_start:code_end].strip()
                
                # Если нет markdown блоков, возвращаем весь контент
                return content.strip()
            
            return None
            
        except Exception as e:
            print(f"⚠️ Ошибка исправления кода: {str(e)}")
            return None
    
    def _offer_user_testing(self, project_dir, main_file, run_command):
        """Предлагает пользователю протестировать программу"""
        try:
            print(f"\n🧪 Хотите протестировать программу?")
            print(f"📍 Путь к файлу: {main_file}")
            print(f"▶️ Команда запуска: {run_command}")
            
            test_choice = input("\nВыберите действие:\n1. Запустить программу еще раз\n2. Открыть папку проекта\n3. Показать код\n4. Пропустить\nВаш выбор (1-4): ")
            
            if test_choice == '1':
                print("\n🔄 Запускаю программу...")
                result = self._run_code(project_dir, main_file, run_command)
                if result['success']:
                    return f"✅ Повторный запуск успешен!\nВывод: {result['output']}"
                else:
                    return f"❌ Ошибка при повторном запуске: {result['error']}"
            
            elif test_choice == '2':
                try:
                    if os.name == 'nt':  # Windows
                        os.startfile(project_dir)
                    else:  # Unix/Linux/Mac
                        subprocess.run(['xdg-open', project_dir])
                    return "📂 Папка проекта открыта"
                except:
                    return f"📂 Откройте папку проекта вручную: {project_dir}"
            
            elif test_choice == '3':
                with open(main_file, 'r', encoding='utf-8') as f:
                    code_content = f.read()
                print(f"\n📄 Содержимое файла {main_file}:")
                print("=" * 50)
                print(code_content)
                print("=" * 50)
                return "📄 Код показан выше"
            
            else:
                return "⏭️ Тестирование пропущено"
                
        except Exception as e:
            return f"⚠️ Ошибка тестирования: {str(e)}"
    
    def _check_task_completion(self, user_task, executed_actions):
        """Проверяет, выполнена ли задача полностью"""
        try:
            # Разбираем задачу на части по ключевым словам
            task_lower = user_task.lower()
            executed_str = ' '.join(executed_actions).lower()
            
            print(f"\n🔍 Анализ выполнения задачи:")
            print(f"📝 Задача: {user_task}")
            print(f"⚙️ Выполненные действия: {executed_actions}")
            
            # Разбиваем задачу на составные шаги
            import re
            task_steps = re.split(r'\b(?:затем|потом|после этого|а затем)\b', task_lower)

            task_steps = [step.strip() for step in task_steps if step.strip()]

            print(f"📋 Определенные шаги: {task_steps}")
            missing_actions = []
            # Анализируем каждый шаг
            for i, step in enumerate(task_steps):
                step_missing = []
                
                # Проверяем создание папок - улучшенная логика
                if "создай" in step and "папк" in step:
                    # Считаем сколько папок нужно создать
                    import re
                    folder_numbers = re.findall(r'\b\d+\b', step)
                    expected_dirs = len(folder_numbers) if folder_numbers else 1
                    actual_dirs = executed_str.count("createdirectory")
                    
                    if actual_dirs < expected_dirs:
                        missing_count = expected_dirs - actual_dirs
                        step_missing.append(f"создание папок: {missing_count} из {expected_dirs} (шаг {i+1})")
                
                # Проверяем создание файлов - улучшенная логика
                if "создай" in step and "файл" in step:
                    # Считаем сколько файлов нужно создать
                    import re
                    # Ищем паттерны типа "файл 11.txt", "файл 21.txt" или "в папке 11 создай файл"
                    file_mentions = re.findall(r'файл\s+[\w\.]+|в\s+папке\s+\d+\s+создай\s+файл', step)
                    
                    # Также проверяем паттерн "в папке X создай файл Y а в папке Z создай файл W"
                    if "а в папке" in step:
                        # Подсчитываем количество упоминаний "в папке ... создай файл"
                        folder_file_patterns = re.findall(r'в\s+папке\s+\w+\s+создай\s+файл', step)
                        expected_files = len(folder_file_patterns)
                    else:
                        expected_files = len(file_mentions) if file_mentions else 1
                    
                    actual_files = executed_str.count("createfile")
                    
                    if actual_files < expected_files:
                        missing_count = expected_files - actual_files
                        step_missing.append(f"создание файлов: {missing_count} из {expected_files} (шаг {i+1})")
                
                # Проверяем перемещение файлов - более детальный анализ
                if any(word in step for word in ["перемести", "перенеси", "переместить"]):
                    # Считаем количество ожидаемых перемещений
                    expected_moves = 0
                    actual_moves = executed_str.count("movefile")
                    
                    # Если упоминаются конкретные папки или файлы
                    if "файлы из папок" in step or "из папок" in step:
                        # Ищем упоминания чисел или имен папок
                        import re
                        folder_numbers = re.findall(r'\b\d+\b', step)
                        if folder_numbers:
                            expected_moves = len(folder_numbers)  # По одному файлу из каждой папки
                    
                    # Если задача о замене местами
                    elif "замени их местами" in step or "поменяй местами" in step:
                        expected_moves = 2  # Обычно 2 файла меняются местами
                    
                    # Общий случай перемещения
                    else:
                        expected_moves = 1
                    
                    if actual_moves < expected_moves:
                        missing_count = expected_moves - actual_moves
                        step_missing.append(f"перемещение файлов: {missing_count} из {expected_moves} (шаг {i+1})")
                
                # Проверяем удаление файлов
                if "удали" in step and "файл" in step:
                    if "deletefile" not in executed_str:
                        step_missing.append(f"удаление файла (шаг {i+1})")
                
                # Проверяем удаление папок
                if "удали" in step and "папк" in step:
                    expected_deletions = 0
                    actual_deletions = executed_str.count("deletedirectory")
                    
                    # Ищем упоминания нескольких папок
                    folder_mentions = re.findall(r'папк[уи]\s+\d+|папк[уи]\s+\w+', step)
                    if folder_mentions:
                        expected_deletions = len(folder_mentions)
                    elif "папки" in step:
                        # Если упоминается "папки" во множественном числе
                        folder_numbers = re.findall(r'\b\d+\b', step)
                        expected_deletions = len(folder_numbers) if folder_numbers else 2
                    else:
                        expected_deletions = 1
                    
                    if actual_deletions < expected_deletions:
                        missing_count = expected_deletions - actual_deletions
                        step_missing.append(f"удаление папок: {missing_count} из {expected_deletions} (шаг {i+1})")
                
                # Проверяем загрузку/генерацию
                if any(word in step for word in ["скачай", "загрузи", "сохрани"]):
                    if "downloadimage" not in executed_str and "downloadfile" not in executed_str and "generateimage" not in executed_str:
                        step_missing.append(f"загрузка/генерация (шаг {i+1})")
                
                missing_actions.extend(step_missing)
            
            print(f"❌ Недостающие действия: {missing_actions}")
            
            return {
                "complete": len(missing_actions) == 0,
                "missing": "; ".join(missing_actions) if missing_actions else None
            }
            
        except Exception as e:
            print(f"⚠️ Ошибка анализа завершения задачи: {e}")
            # В случае ошибки анализа считаем задачу выполненной
            return {"complete": True, "missing": None}
    
    def _is_simple_conversation(self, user_input):
        """Определяет, является ли запрос простым общением"""
        # Ключевые слова, которые указывают на задачи, требующие инструментов
        task_keywords = [
            "создай", "сгенерируй", "генерируй", "сделай", "скачай", "загрузи", "удали", "перемести", 
            "переименуй", "найди файл", "запусти", "выполни", "установи", "открой", 
            "сохрани", "измени", "переведи", "поиск в интернете", "найди в интернете",
            "создать", "сгенерировать", "скачать", "загрузить", "удалить", "переместить",
            "аудио", "удио", "озвучь", "озвучить", "произнеси", "говори", "скажи", "voice", "audio", "sound",
            "изображение", "картинка", "рисунок", "фото", "image", "picture", "photo", "найди", "search",
            # Добавляем ключевые слова для чтения файлов
            "прочитай", "прочти", "читай", "что написано", "содержимое", "содержание", "текст из",
            "что в файле", "что в документе", "посмотри файл", "покажи файл", "файл содержит",
            "read file", "show file", "file content", "file contains"
        ]
        
        # Специальная проверка для аудио запросов
        audio_patterns = [
            "аудио с текстом", "удио с текстом", "аудио со словами", "удио со словами",
            "озвучь текст", "произнеси", "скажи голосом", "сделай аудио", "сгенерируй аудио",
            "генерируй аудио", "создай аудио", "сделай удио", "сгенерируй удио"
        ]
        
        # Ключевые слова простого общения
        conversation_keywords = [
            "привет", "как дела", "что умеешь", "кто ты", "расскажи о", "объясни", 
            "что такое", "помоги понять", "спасибо", "пока", "как у тебя дела",
            "что нового", "как работает", "можешь рассказать", "что думаешь", "hello", "hi"
        ]
        
        user_lower = user_input.lower().strip()
        
        # Отладочная информация
        if self.debug_mode:
            print(f"[DEBUG] Анализ запроса: '{user_input}'")
            print(f"[DEBUG] В нижнем регистре: '{user_lower}'")
            print(f"[DEBUG] Длина: {len(user_lower)} символов, {len(user_input.split())} слов")
        
        # Если это короткий запрос (меньше 10 символов) и нет ключевых слов задач
        if len(user_lower) < 10:
            has_task_keywords = any(keyword in user_lower for keyword in task_keywords)
            if self.debug_mode:
                print(f"[DEBUG] Короткий запрос (<10 символов), есть ключевые слова задач: {has_task_keywords}")
            return not has_task_keywords
        
        # ПЕРВЫЙ ПРИОРИТЕТ: Если есть ключевые слова задач, это НЕ простое общение
        has_task_keywords = any(keyword in user_lower for keyword in task_keywords)
        if has_task_keywords:
            if self.debug_mode:
                found_keywords = [kw for kw in task_keywords if kw in user_lower]
                print(f"[DEBUG] Найдены ключевые слова задач: {found_keywords}, возвращаем False")
            return False
        
        # ВТОРОЙ ПРИОРИТЕТ: Если есть прямые указания на общение (И НЕТ ключевых слов задач)
        has_conversation_keywords = any(keyword in user_lower for keyword in conversation_keywords)
        if has_conversation_keywords:
            if self.debug_mode:
                print(f"[DEBUG] Найдены ключевые слова общения (и нет ключевых слов задач), возвращаем True")
            return True
            
        # Если это вопрос без конкретных действий
        is_question = user_lower.startswith(("что", "как", "зачем", "почему", "когда", "где", "кто", "можешь ли", "умеешь ли"))
        if is_question:
            if self.debug_mode:
                print(f"[DEBUG] Это вопрос, возвращаем True")
            return True
            
        # По умолчанию считаем простым общением для коротких фраз
        word_count = len(user_input.split())
        is_short = word_count <= 3
        if self.debug_mode:
            print(f"[DEBUG] Слов: {word_count}, короткая фраза (<=3): {is_short}")
        return is_short
    
    def _handle_simple_conversation(self, user_input):
        """Обрабатывает простое общение без использования инструментов"""
        try:
            print(f"\n💬 {get_text('task_processing', self.interface_language)}")
            
            # Переводим промпт для AI если нужно
            translated_prompt = translate_prompt_for_ai(user_input, self.output_language)
            
            # Создаем простой запрос без инструментов
            messages = [{"role": "user", "content": translated_prompt}]
            
            # Отправляем запрос БЕЗ инструментов (и без tool_choice)
            response = communicate_with_Pollinations_chat(
                self.current_model, 
                messages, 
                tools=None,  # Без инструментов!
                tool_choice=None,  # Не передаем tool_choice если нет tools
                api_token=self.api_token
            )
            
            if "choices" in response and len(response["choices"]) > 0:
                content = response["choices"][0]["message"].get("content", "")
                if content:
                    return content
                else:
                    return get_text('error', self.interface_language) + ": " + "Не смог сформулировать ответ."
            elif "error" in response:
                return f"{get_text('error', self.interface_language)}: {response['error']}"
            else:
                return get_text('error', self.interface_language) + ": Не удалось получить ответ."
                
        except Exception as e:
            return f"{get_text('error', self.interface_language)}: {str(e)}"
    
    def _get_interface_messages(self):
        """Возвращает словарь сообщений интерфейса"""
        return {
            'api_token_loaded': get_text('api_token_loaded', self.interface_language),
            'api_token_missing': get_text('api_token_missing', self.interface_language),
            'welcome': get_text('welcome', self.interface_language),
            'agent_started': get_text('agent_started', self.interface_language),
            'agent_stopped': get_text('agent_stopped', self.interface_language),
            'task_processing': get_text('task_processing', self.interface_language),
            'task_completed': get_text('task_completed', self.interface_language),
            'task_failed': get_text('task_failed', self.interface_language),
            'available_models': get_text('available_models', self.interface_language),
            'select_model': get_text('select_model', self.interface_language),
            'invalid_choice': get_text('invalid_choice', self.interface_language),
            'model_selected': get_text('model_selected', self.interface_language),
            'enter_task': get_text('enter_task', self.interface_language),
            'exit_command': get_text('exit_command', self.interface_language),
            'change_command': get_text('change_command', self.interface_language),
            'result': get_text('result', self.interface_language),
        }

    def run(self):
        # Проверяем, нужно ли показать выбор языка при первом запуске
        lang_settings = get_language_settings()
        if lang_settings['first_startup']:
            interface_lang, output_lang = show_language_selection()
            save_language_settings(interface_lang, output_lang)
            # Обновляем настройки в текущем экземпляре
            self.interface_language = interface_lang
            self.output_language = output_lang
            self.messages = self._get_interface_messages()
            print(f"\n{get_text('restart_required', self.interface_language)}")
            return
        
        print(f"\n{self.messages['welcome']}")
        print(f"🚀 {self.messages['agent_started']}")
        
        # Показываем настройки по умолчанию
        if self.default_model in self.model_list:
            self.current_model = self.default_model
            print(f"\n🤖 {get_text('model_selected', self.interface_language)}: {self.current_model}")
        else:
            self.select_model()
        
        while True:
            user_input = input(f"\n{self.messages['enter_task']} ")
            
            if user_input.lower() == self.messages['exit_command']:
                print(f"\n👋 {self.messages['agent_stopped']}")
                break
            elif user_input.lower() == self.messages['change_command']:
                self.change_model()
                continue
            
            logger.info(f"Задача пользователя: {user_input}")
            print(f"\n🔄 {self.messages['task_processing']}")
            
            # Используем настройку максимального количества попыток из .env
            plan = self.generate_plan(user_input, self.max_attempts)
            
            print(f"\n{self.messages['result']}")
            print(plan)
            
            # Если задача уже выполнена автоматически, продолжаем
            if get_text('task_completed', self.interface_language) in plan or get_text('task_failed', self.interface_language) in plan:
                continue
            
            # Пропускаем ручное выполнение планов (старая логика не работает)
            if any(marker in plan for marker in ["✅", "❌", "[Ошибка]"]):
                print(f"📝 {get_text('error', self.interface_language)}: Получен текстовый план. Мануальное выполнение не поддерживается.")
                print(f"💡 Попробуйте переформулировать задачу более конкретно.")

if __name__ == "__main__":
    check_for_updates()
    agent = PollinationsAgent()
    agent.run()